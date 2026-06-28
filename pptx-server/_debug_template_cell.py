# -*- coding: utf-8 -*-
"""양식 PPTX에서 공휴일 셀(1주차 수요일=col3) 서식 확인"""
import sys, copy
sys.stdout.reconfigure(encoding='utf-8')
from lxml import etree
from pptx import Presentation

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

TPL = 'templates/2026년 6월 식단표(양식).pptx'
prs = Presentation(TPL)

# 슬라이드 0, 테이블 찾기
slide = list(prs.slides)[0]
table = None
for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        break

if table is None:
    print('테이블 없음')
    sys.exit()

print(f'테이블: {len(table.rows)}행 × {len(table.columns)}열')
print()

# 1주차 선거일 = row0(lunch)/row1(am)/row2(pm), col=3 (0-indexed)
# col0=라벨, col1~5=월~금 → 수요일=col3
TARGET_COL = 3

for ri in [0, 1, 2]:
    sec = ['lunch', 'am', 'pm'][ri]
    cell = table.cell(ri, TARGET_COL)
    tc   = cell._tc
    txBody = tc.find(f'{{{NS_A}}}txBody')

    # ── 텍스트 ──────────────────────────────────────
    texts = []
    if txBody is not None:
        for p in txBody.findall(f'{{{NS_A}}}p'):
            para = []
            for r in p.findall(f'{{{_NS_A}}}r' if False else f'{{{NS_A}}}r'):
                t = r.find(f'{{{NS_A}}}t')
                if t is not None and t.text:
                    para.append(t.text)
            texts.append(''.join(para))
    text_repr = ' / '.join(f'<{t}>' for t in texts)

    # ── 배경색 (tcPr > solidFill) ──────────────────
    tcPr = tc.find(f'{{{NS_A}}}tcPr')
    bg_color = None
    if tcPr is not None:
        solidFill = tcPr.find(f'.//{{{NS_A}}}solidFill')
        if solidFill is not None:
            srgb = solidFill.find(f'{{{NS_A}}}srgbClr')
            sysClr = solidFill.find(f'{{{NS_A}}}sysClr')
            if srgb is not None:
                bg_color = f'srgb #{srgb.get("val","?")}'
            elif sysClr is not None:
                bg_color = f'sysClr {sysClr.get("lastClr","?")}'

    # ── 음영/패턴 (prstClr, gradFill, pattFill 등) ─
    extra_fill = None
    if tcPr is not None:
        for tag in ('gradFill', 'pattFill', 'blipFill', 'noFill'):
            el = tcPr.find(f'{{{NS_A}}}{tag}')
            if el is not None:
                extra_fill = tag
                break

    print(f'row{ri:02d}({sec:6s}) col{TARGET_COL}')
    print(f'  텍스트: {text_repr if text_repr else "(비어있음)"}')
    print(f'  배경색: {bg_color if bg_color else "(없음/투명)"}')
    print(f'  특수fill: {extra_fill if extra_fill else "(없음)"}')

    # ── tcPr XML 전체 ───────────────────────────────
    if tcPr is not None:
        xml_str = etree.tostring(tcPr, pretty_print=True).decode('utf-8')
        # 길면 잘라서 보여줌
        lines = xml_str.strip().splitlines()
        if len(lines) > 20:
            lines = lines[:20] + ['  ... (생략)']
        print(f'  tcPr XML:')
        for ln in lines:
            print(f'    {ln}')
    else:
        print('  tcPr: 없음')
    print()
