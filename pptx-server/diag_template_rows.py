"""
진단 1: 양식.pptx 원본 간식행 높이 확인
진단 2: 간식 셀 <a:bodyPr> autofit 설정 확인
"""
import os, sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from lxml import etree

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

TEMPLATE = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        'templates', '2026년 6월 식단표(양식).pptx')

prs = Presentation(TEMPLATE)

for si, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        tbl = shape.table
        n = len(tbl.rows)
        print(f'\n[slide{si+1}]  행 수={n}')

        for ri, row in enumerate(tbl.rows):
            cell_text = tbl.cell(ri, 0).text.strip().replace('\n', ' ')[:20]
            print(f'  행[{ri:02d}] h={row.height or 0:>9,}  "{cell_text}"')

        # ── bodyPr 확인: 간식행 첫 번째 셀 ──────────────────────────
        # 간식행 = height가 가장 작은 행들
        min_h = min(r.height or 0 for r in tbl.rows)
        print(f'\n  간식행 기준 높이(최솟값): {min_h:,}')

        # 첫 번째 간식행의 첫 번째 셀 bodyPr 덤프
        for ri, row in enumerate(tbl.rows):
            if (row.height or 0) == min_h:
                cell = tbl.cell(ri, 0)
                tc_el = cell._tc
                txBody = tc_el.find(f'{{{NS_A}}}txBody')
                if txBody is not None:
                    bodyPr = txBody.find(f'{{{NS_A}}}bodyPr')
                    if bodyPr is not None:
                        attrs = dict(bodyPr.attrib)
                        # spAutoFit / normAutofit / noAutofit 자식 요소
                        children = [etree.QName(c).localname for c in bodyPr]
                        print(f'  bodyPr attrs  : {attrs}')
                        print(f'  bodyPr children: {children}')
                    else:
                        print('  bodyPr: 없음')
                break   # 첫 번째 간식행만
        break           # 슬라이드별 첫 번째 테이블만
