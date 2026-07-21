from pptx import Presentation

prs = Presentation("templates/2026년 6월 식단표(양식).pptx")
slide = prs.slides[0]

for shape in slide.shapes:
    if shape.name in ["ORIGIN_BOX", "MATERIAL_BOX", "ALLERGY_BOX", "그림 224", "그래픽 118"]:
        print(f"{shape.name}: top={shape.top}, left={shape.left}, w={shape.width}, h={shape.height}")
    if shape.has_table:
        table = shape.table
        print(f"\n테이블: {len(table.rows)}행")
        print(f"테이블 top={shape.top}")
        for i, row in enumerate(table.rows):
            print(f"  행[{i}] height={row.height}")
