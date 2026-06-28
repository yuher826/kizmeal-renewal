# -*- coding: utf-8 -*-
"""
z-order 정밀 확인 + 그림 4 / TextBox 26 분석
- spTree 직접 자식을 document order로 순회 → 정확한 z-order
- 그림 4가 표 앞/뒤인지 확정
- 템플릿에도 같은 도형이 있는지 확인
"""
import sys, zipfile
sys.stdout.reconfigure(encoding='utf-8')
from lxml import etree
from pptx import Presentation

NS_A   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_PML = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_R   = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

# ── 3일 셀 절대 좌표 (이미 계산된 값) ────────────────────────────────
CELL_X1, CELL_Y1 = 3_322_313, 1_156_203
CELL_X2, CELL_Y2 = 4_668_713, 2_349_546

def overlaps(ex, ey, ecx, ecy):
    ex2, ey2 = ex+ecx, ey+ecy
    return not (ex2 <= CELL_X1 or CELL_X2 <= ex or ey2 <= CELL_Y1 or CELL_Y2 <= ey)

def get_xfrm(el):
    # 직접 자식에서 xfrm 검색 (sp, graphicFrame 등)
    for ns in (NS_PML, NS_A):
        xfrm = el.find(f'{{{ns}}}xfrm')
        if xfrm is not None:
            off = xfrm.find(f'{{{NS_A}}}off')
            ext = xfrm.find(f'{{{NS_A}}}ext')
            if off is not None and ext is not None:
                return (int(off.get('x',0)), int(off.get('y',0)),
                        int(ext.get('cx',0)), int(ext.get('cy',0)))
    # pic 등: spPr → a:xfrm 안에 위치
    for spPr in el.findall(f'.//{{{NS_PML}}}spPr') + el.findall(f'.//{{{NS_A}}}spPr'):
        xfrm = spPr.find(f'{{{NS_A}}}xfrm')
        if xfrm is not None:
            off = xfrm.find(f'{{{NS_A}}}off')
            ext = xfrm.find(f'{{{NS_A}}}ext')
            if off is not None and ext is not None:
                return (int(off.get('x',0)), int(off.get('y',0)),
                        int(ext.get('cx',0)), int(ext.get('cy',0)))
    return None

def get_name(el):
    cNvPr = el.find(f'.//{{{NS_PML}}}cNvPr')
    return cNvPr.get('name','?') if cNvPr is not None else '?'

# ════════════════════════════════════════════════════════════════════
# 1. spTree 직접 자식만 document order로 z-order 확인
# ════════════════════════════════════════════════════════════════════
print('=' * 65)
print('강동E output — spTree 직접 자식 z-order (뒤→앞 순)')
print('3일 셀 영역: x1=3,322,313 y1=1,156,203 x2=4,668,713 y2=2,349,546')
print('=' * 65)

src = 'output/강동E_202606.pptx'
prs = Presentation(src)
slide = list(prs.slides)[0]
spTree = slide.shapes._spTree     # p:spTree가 도형 컨테이너

direct_children = list(spTree)   # 직접 자식만 (중첩 불포함)
table_zidx = None

for zi, el in enumerate(direct_children):
    tag  = el.tag.split('}')[-1]
    name = get_name(el)
    rect = get_xfrm(el)
    if rect:
        ex, ey, ecx, ecy = rect
        hit = overlaps(ex, ey, ecx, ecy)
        flag = '  ⚠️ 겹침!' if hit else ''
        print(f'  z={zi:3d} [{tag:14s}] {name:25s} x={ex:>10,} y={ey:>10,}'
              f' cx={ecx:>9,} cy={ecy:>9,}{flag}')
        if tag == 'graphicFrame' and name == 'MENU_TABLE':
            table_zidx = zi
    else:
        print(f'  z={zi:3d} [{tag:14s}] {name:25s} (xfrm 없음)')

print()
if table_zidx is not None:
    print(f'MENU_TABLE z-index = {table_zidx}')
    print(f'→ z-index > {table_zidx} 인 도형이 표 앞(표 위)에 있음')


# ════════════════════════════════════════════════════════════════════
# 2. 그림 4 상세 분석 (실제 이미지 파일, 관계 추적)
# ════════════════════════════════════════════════════════════════════
print()
print('=' * 65)
print('그림 4 상세 분석 — 이미지 관계(rId) 추적')
print('=' * 65)

with zipfile.ZipFile(src, 'r') as zf:
    # slide1.xml.rels 에서 rId→이미지 파일 매핑
    rels_xml = zf.read('ppt/slides/_rels/slide1.xml.rels')
    rels_root = etree.fromstring(rels_xml)
    rid_map = {}
    for rel in rels_root:
        rid_map[rel.get('Id')] = rel.get('Target', '')

    # spTree에서 그림 4 찾기 (직접 자식 + 중첩 포함)
    for el in spTree.iter():
        name = get_name(el)
        if name == '그림 4':
            tag = el.tag.split('}')[-1]
            rect = get_xfrm(el)
            print(f'  태그: {tag}, 이름: {name}')
            if rect:
                ex, ey, ecx, ecy = rect
                print(f'  좌표: x={ex:,} y={ey:,} cx={ecx:,} cy={ecy:,}')
                print(f'  크기: {ecx/914400*2.54:.2f}cm × {ecy/914400*2.54:.2f}cm')

            # 이미지 rId 추출 (blip rEmbed)
            blip = el.find(f'.//{{{NS_A}}}blip')
            if blip is not None:
                rId = blip.get(f'{{{NS_R}}}embed')
                target = rid_map.get(rId, '(알 수 없음)')
                print(f'  이미지 rId: {rId} → {target}')
                # 이미지 파일 존재 확인
                img_path = 'ppt/slides/' + target if target.startswith('../') else \
                           'ppt/' + target.lstrip('../')
                img_path = img_path.replace('ppt/slides/../', 'ppt/')
                exists = img_path in zf.namelist()
                print(f'  이미지 파일 경로: {img_path} (존재={exists})')
                if exists:
                    size = len(zf.read(img_path))
                    print(f'  이미지 파일 크기: {size:,} bytes')
            else:
                print('  이미지 blip 없음 (그림이 아닌 도형일 수 있음)')
            # XML 일부
            xml_str = etree.tostring(el, pretty_print=True).decode('utf-8')
            lines = xml_str.strip().splitlines()[:30]
            print('  XML (앞 30줄):')
            for ln in lines:
                print(f'    {ln}')
            break


# ════════════════════════════════════════════════════════════════════
# 3. 템플릿 원본에도 그림 4가 있는지 확인
# ════════════════════════════════════════════════════════════════════
print()
print('=' * 65)
print('템플릿 원본(양식.pptx) — S1(15행) 슬라이드 그림 4 확인')
print('=' * 65)

TPL = 'templates/2026년 6월 식단표(양식).pptx'
prs_t = Presentation(TPL)

for si, slide_t in enumerate(prs_t.slides):
    print(f'\n── 슬라이드 {si} (S{si+1}) ──')
    spTree_t = slide_t.shapes._spTree

    # 이름으로 그림 4 직접 검색
    pic4_found = False
    for el in spTree_t.iter():
        name = get_name(el)
        if name == '그림 4':
            tag  = el.tag.split('}')[-1]
            rect = get_xfrm(el)
            pic4_found = True
            if rect:
                ex, ey, ecx, ecy = rect
                hit = overlaps(ex, ey, ecx, ecy)
                print(f'  [그림 4 발견] [{tag}] x={ex:,} y={ey:,} cx={ecx:,} cy={ecy:,}'
                      f'  {"⚠️ 3일 겹침" if hit else ""}')
            else:
                print(f'  [그림 4 발견] [{tag}] xfrm 없음')
    if not pic4_found:
        print('  그림 4 없음')

    # 겹치는 도형 전체 (고침된 get_xfrm 사용)
    overlapping = []
    for el in spTree_t.iter():
        rect = get_xfrm(el)
        if rect:
            ex, ey, ecx, ecy = rect
            if overlaps(ex, ey, ecx, ecy):
                name = get_name(el)
                tag  = el.tag.split('}')[-1]
                overlapping.append((tag, name, ex, ey, ecx, ecy))
    if overlapping:
        print('  3일 셀 영역 겹치는 도형:')
        for tag, name, ex, ey, ecx, ecy in overlapping:
            print(f'    [{tag}] {name}  x={ex:,} y={ey:,} cx={ecx:,} cy={ecy:,}')
