# -*- coding: utf-8 -*-
"""
키즈밀 식단표 PPTX 생성기 (빈 양식 1개 → 49개원 자동생성)
========================================================================

★ 대원칙: 빈 양식 슬라이드 XML을 zipfile 수준에서 재조합 후
  <a:t> 텍스트만 교체 (cell.text= 절대 금지).

빈 양식 슬라이드 구성:
  S1 (idx=0): 15행 = (중식+오전+오후) × 5주   (타입C/D)
  S2 (idx=1): 10행 = (중식+오전) × 5주         (타입A/B/E/F/G/songpaE)
  S3 (idx=2):  5행 = 오후간식 × 5주            (타입B/F/G)

셀 XML 구조 (양식.pptx 기준):
  중식 셀: <a:p>... <a:r>행1</a:r><a:br/><a:r>행2</a:r>... (7 run + 6 br)
  간식 셀: <a:p>... <a:r>메뉴</a:r>  (1 run; kcal은 <a:br>+<a:r> 추가)
  빈 셀(col2-5, 2주~5주): run 없음 → _init_table_cells로 구조 복사

⚠️  일본어 금지. 모든 주석/문자열 한국어.
"""

import copy
import io
import re
import zipfile

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from bracket_parser import resolve_for_branch

# ── 알러지 원문자 범위 ①~⑲ ─────────────────────────────────────────
_ALLERGY_LO, _ALLERGY_HI = 0x2460, 0x2472

# ── XML 네임스페이스 ───────────────────────────────────────────────
_NS_A   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_NS_PML = 'http://schemas.openxmlformats.org/presentationml/2006/main'
_NS_R   = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
_NS_PKG = 'http://schemas.openxmlformats.org/package/2006/relationships'
_NS_CT  = 'http://schemas.openxmlformats.org/package/2006/content-types'

_SLIDE_REL_TYPE = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide'

_S1, _S2, _S3 = 0, 1, 2

# ── 공휴일에도 운영하는 원 ────────────────────────────────────────
_HOLIDAY_OPERATING = {'덕양P', '광교SLP'}

# ── 날짜·사유 텍스트 색 (원본 정답본 #318E72 초록으로 통일) ─────────
# 평일/공휴일/선거일 날짜 숫자 + 공휴일 사유 텍스트 모두 이 색 사용.
# 알레르기 번호색(#FF0000)·기타 박스색과는 별개 — 충돌 없음.
DATE_COLOR_GREEN = '318E72'

# ════════════════════════════════════════════════════════════════════
# 1. 타입별 슬라이드 플랜
# ════════════════════════════════════════════════════════════════════
TYPE_PLANS = {
    'A':       [(_S2, ['lunch', 'am'],            {})],
    'B':       [(_S2, ['lunch', 'am'],            {}),
                (_S3, ['pm'],                     {})],
    'C':       [(_S1, ['lunch', 'am', 'pm'],      {})],
    'D':       [(_S1, ['lunch', 'am', 'pm'],      {}),
                (_S1, ['lunch', 'am', 'pm'],      {})],
    'E':       [(_S2, ['lunch', 'am'],            {}),
                (_S1, [],                          {'skip': True}),
                (_S3, [],                          {'skip': True})],
    'F':       [(_S2, ['lunch', 'am'],            {}),
                (_S2, ['lunch', 'am'],             {'eng': True}),
                (_S3, ['pm'],                      {}),
                (_S3, ['pm'],                      {'eng': True})],
    'G':       [(_S2, ['lunch', 'am'],            {}),
                (_S3, ['pm'],                      {}),
                (_S2, ['lunch', 'am'],             {'strip': True}),
                (_S3, ['pm'],                      {'strip': True})],
    'songpaE': [(_S2, ['lunch', 'am'],            {}),
                (_S2, ['pm', 'care'],              {'care_label': '저녁'})],
}


# ════════════════════════════════════════════════════════════════════
# 2. ZIP 기반 슬라이드 플랜 → Presentation
# ════════════════════════════════════════════════════════════════════
def build_pptx_from_plan(template_path, slide_indices):
    with zipfile.ZipFile(template_path, 'r') as zf:
        file_map = {n: zf.read(n) for n in zf.namelist()}

    prs_xml  = etree.fromstring(file_map['ppt/presentation.xml'])
    sldIdLst = prs_xml.find(f'.//{{{_NS_PML}}}sldIdLst')
    prs_rels = etree.fromstring(file_map['ppt/_rels/presentation.xml.rels'])

    rId_to_target = {}
    for rel in prs_rels:
        t = rel.get('Type', '')
        if ('slide' in t and 'slideLayout' not in t
                and 'slideMaster' not in t and 'notes' not in t.lower()):
            rId_to_target[rel.get('Id')] = rel.get('Target')

    template_slide_files = []
    for sldId in sldIdLst:
        rId = sldId.get(f'{{{_NS_R}}}id')
        target = rId_to_target.get(rId)
        if target:
            template_slide_files.append('ppt/' + target)

    new_files = {}
    for name, data in file_map.items():
        is_slide = name.startswith('ppt/slides/slide') and name.endswith('.xml')
        is_notes = name.startswith('ppt/notesSlides/')
        if not is_slide and not is_notes:
            new_files[name] = data

    new_targets = []
    for i, t_idx in enumerate(slide_indices):
        src = template_slide_files[t_idx]
        dst = f'ppt/slides/slide{i+1}.xml'
        new_files[dst] = file_map[src]

        src_rels = src.replace('/slides/', '/slides/_rels/').replace('.xml', '.xml.rels')
        dst_rels = f'ppt/slides/_rels/slide{i+1}.xml.rels'
        if src_rels in file_map:
            rx = etree.fromstring(file_map[src_rels])
            for rel in list(rx):
                if 'notes' in rel.get('Type', '').lower():
                    rx.remove(rel)
            new_files[dst_rels] = etree.tostring(
                rx, xml_declaration=True, encoding='UTF-8', standalone=True)

        new_targets.append(f'slides/slide{i+1}.xml')

    for child in list(sldIdLst):
        sldIdLst.remove(child)
    for rel in list(prs_rels):
        t = rel.get('Type', '')
        if ('slide' in t and 'slideLayout' not in t
                and 'slideMaster' not in t and 'notes' not in t.lower()):
            prs_rels.remove(rel)

    for i, target in enumerate(new_targets):
        rid = f'rId_s{i+1}'
        etree.SubElement(prs_rels, f'{{{_NS_PKG}}}Relationship', {
            'Id': rid, 'Type': _SLIDE_REL_TYPE, 'Target': target,
        })
        elem = etree.SubElement(sldIdLst, f'{{{_NS_PML}}}sldId', {'id': str(300 + i)})
        elem.set(f'{{{_NS_R}}}id', rid)

    ct_xml   = etree.fromstring(new_files['[Content_Types].xml'])
    slide_ct = 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml'
    for ov in list(ct_xml):
        pn = ov.get('PartName', '')
        if '/slides/slide' in pn and 'Layout' not in pn and 'Master' not in pn:
            ct_xml.remove(ov)
    for i in range(len(new_targets)):
        etree.SubElement(ct_xml, f'{{{_NS_CT}}}Override', {
            'PartName': f'/ppt/slides/slide{i+1}.xml',
            'ContentType': slide_ct,
        })

    def _b(elem):
        return etree.tostring(elem, xml_declaration=True, encoding='UTF-8', standalone=True)

    new_files['ppt/presentation.xml']           = _b(prs_xml)
    new_files['ppt/_rels/presentation.xml.rels'] = _b(prs_rels)
    new_files['[Content_Types].xml']             = _b(ct_xml)

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        for name, data in new_files.items():
            zf.writestr(name, data)
    buf.seek(0)
    return Presentation(buf)


# ════════════════════════════════════════════════════════════════════
# 3. 표 셀 XML 초기화 (빈 셀 + 타입 불일치 셀 → 올바른 구조로 교체)
# ════════════════════════════════════════════════════════════════════
def _normalize_snack_tmpl(txBody):
    """간식 txBody를 1단락·1run·빈텍스트로 정규화 (오염 템플릿 대비)."""
    all_ps = txBody.findall(f'{{{_NS_A}}}p')
    for extra_p in all_ps[1:]:
        txBody.remove(extra_p)
    p = txBody.find(f'{{{_NS_A}}}p')
    if p is None:
        return txBody
    runs = p.findall(f'{{{_NS_A}}}r')
    brs  = p.findall(f'{{{_NS_A}}}br')
    for r in runs[1:]:
        p.remove(r)
    for br in brs:
        p.remove(br)
    if runs:
        t = runs[0].find(f'{{{_NS_A}}}t')
        if t is not None:
            t.text = ''
    return txBody


def _make_snack_tmpl_from_lunch(lunch_tmpl):
    """중식 txBody에서 1-run 간식 txBody 파생 (snack-only 슬라이드 폴백)."""
    tmpl = copy.deepcopy(lunch_tmpl)
    # 여분 단락 제거 (다단락 lunch_tmpl이면 중식 데이터가 snack에 남음)
    all_ps = tmpl.findall(f'{{{_NS_A}}}p')
    for extra_p in all_ps[1:]:
        tmpl.remove(extra_p)
    p = tmpl.find(f'{{{_NS_A}}}p')
    if p is None:
        return tmpl
    runs = p.findall(f'{{{_NS_A}}}r')
    brs  = p.findall(f'{{{_NS_A}}}br')
    for r in runs[1:]:
        p.remove(r)
    for br in brs:
        p.remove(br)
    if runs:
        t = runs[0].find(f'{{{_NS_A}}}t')
        if t is not None:
            t.text = ''
    return tmpl


def _init_table_cells(table, sections):
    """
    데이터 셀(col1~N)의 txBody를 섹션 타입에 맞게 초기화.
    - run 없는 빈 셀 → 적절한 template 복사
    - 타입 불일치(lunch template in snack row) → 올바른 template으로 교체
    """
    n_cols = len(table.columns)
    n_secs = len(sections)

    # sections 기반으로 lunch/snack template 추출 (run 수 임계값 불사용).
    # ★col1(월요일) 고정이 아니라 데이터 열 전체(col1~n_cols-1)를 순회한다★
    # — 1일이 월요일이 아닌 달은 디자이너 양식의 월요일 칸이 비어 있어
    # (9월 실측: 슬라이드1 중식 행 전부 col1 run수=0, 화요일인 col2에서야
    # run이 나옴) col1만 보면 lunch_tmpl이 끝내 None이 되어 전 주차 중식이
    # 초기화되지 않고 빈칸으로 남는다(HANDOFF "1일이 월요일이 아닌 달"
    # 버그 참고). lunch_tmpl과 snack_tmpl은 각각 독립적으로 찾는다 —
    # 어느 요일에 데이터가 있든 동작해야 한다.
    lunch_tmpl = None
    snack_tmpl = None
    for ri in range(len(table.rows)):
        sec = sections[ri % n_secs]
        if sec == 'lunch' and lunch_tmpl is not None:
            continue
        if sec != 'lunch' and snack_tmpl is not None:
            continue
        for ci in range(1, n_cols):
            src_txBody = table.cell(ri, ci)._tc.find(f'{{{_NS_A}}}txBody')
            if src_txBody is None:
                continue
            n = len(src_txBody.findall(f'.//{{{_NS_A}}}r'))
            if n == 0:
                continue
            if sec == 'lunch':
                lunch_tmpl = copy.deepcopy(src_txBody)
            else:
                snack_tmpl = _normalize_snack_tmpl(copy.deepcopy(src_txBody))
            break
        if lunch_tmpl is not None and snack_tmpl is not None:
            break

    # S3 처럼 snack_tmpl이 없으면 lunch_tmpl 에서 파생
    if snack_tmpl is None and lunch_tmpl is not None:
        snack_tmpl = _make_snack_tmpl_from_lunch(lunch_tmpl)

    def _tmpl(sec):
        return lunch_tmpl if sec == 'lunch' else snack_tmpl

    for ri in range(len(table.rows)):
        sec  = sections[ri % n_secs]
        tmpl = _tmpl(sec)
        if tmpl is None:
            continue
        need_lunch = (sec == 'lunch')
        for ci in range(1, n_cols):
            dst_tc     = table.cell(ri, ci)._tc
            dst_txBody = dst_tc.find(f'{{{_NS_A}}}txBody')
            if dst_txBody is None:
                continue
            cur_n  = len(dst_txBody.findall(f'.//{{{_NS_A}}}r'))
            has_lunch = (cur_n >= 5)
            # 비어있거나 타입 불일치이면 교체
            if cur_n == 0 or (need_lunch != has_lunch):
                dst_tc.replace(dst_txBody, copy.deepcopy(tmpl))


# ════════════════════════════════════════════════════════════════════
# 4. 알러지 유틸
# ════════════════════════════════════════════════════════════════════
def split_menu_allergy(text):
    allergy = ''.join(c for c in text if _ALLERGY_LO <= ord(c) <= _ALLERGY_HI)
    menu    = ''.join(c for c in text if not (_ALLERGY_LO <= ord(c) <= _ALLERGY_HI))
    return menu, allergy


def strip_allergy(text):
    return ''.join(c for c in text if not (_ALLERGY_LO <= ord(c) <= _ALLERGY_HI))


# ════════════════════════════════════════════════════════════════════
# 5. 셀 텍스트 교체 (XML 직접 조작)
# ════════════════════════════════════════════════════════════════════
def _get_p_and_runs(cell):
    """셀 첫 단락(p)과 <a:r> 리스트 반환."""
    tc     = cell._tc
    txBody = tc.find(f'{{{_NS_A}}}txBody')
    if txBody is None:
        return None, []
    p = txBody.find(f'{{{_NS_A}}}p')
    if p is None:
        return None, []
    runs = p.findall(f'{{{_NS_A}}}r')
    return p, runs


def _set_run_text(run, text):
    t_elem = run.find(f'{{{_NS_A}}}t')
    if t_elem is not None:
        t_elem.text = text


def _make_allergy_run(template_run, allergy_text):
    """알레르기 번호 전용 빨간색(FF0000) run 생성 후 반환"""
    new_run = copy.deepcopy(template_run)
    rPr = new_run.find(f'{{{_NS_A}}}rPr')
    if rPr is not None:
        # 기존 solidFill 제거 (list()로 복사본 순회해야 안전)
        for child in list(rPr):
            if child.tag == f'{{{_NS_A}}}solidFill':
                rPr.remove(child)
        # FF0000 solidFill 첫 번째 자식으로 삽입
        solid = etree.Element(f'{{{_NS_A}}}solidFill')
        srgb  = etree.SubElement(solid, f'{{{_NS_A}}}srgbClr')
        srgb.set('val', 'FF0000')
        rPr.insert(0, solid)
    # 텍스트 교체
    t = new_run.find(f'{{{_NS_A}}}t')
    if t is None:
        t = etree.SubElement(new_run, f'{{{_NS_A}}}t')
    t.text = allergy_text
    return new_run


def _apply_inline_allergy(template_run, text, p):
    """간식 셀용: 텍스트 내 알레르기 번호를 인라인 분리해 빨간색 run 삽입.
    예) '삶은계란①, 요구르트②'
        → run(삶은계란) + run_red(①) + run(, 요구르트) + run_red(②)
    """
    segments = []
    current = ''
    is_allergy = None
    for char in text:
        char_is_allergy = _ALLERGY_LO <= ord(char) <= _ALLERGY_HI
        if is_allergy is None:
            is_allergy = char_is_allergy
        if char_is_allergy != is_allergy:
            if current:
                segments.append((current, is_allergy))
            current = char
            is_allergy = char_is_allergy
        else:
            current += char
    if current:
        segments.append((current, is_allergy))

    if not segments:
        _set_run_text(template_run, '')
        return

    first_text, first_is_allergy = segments[0]
    if first_is_allergy:
        _set_run_text(template_run, '')
        insert_pos = list(p).index(template_run) + 1
        p.insert(insert_pos, _make_allergy_run(template_run, first_text))
        insert_pos += 1
    else:
        _set_run_text(template_run, first_text)
        insert_pos = list(p).index(template_run) + 1

    for seg_text, seg_is_allergy in segments[1:]:
        if seg_is_allergy:
            new_run = _make_allergy_run(template_run, seg_text)
        else:
            new_run = copy.deepcopy(template_run)
            t_elem = new_run.find(f'{{{_NS_A}}}t')
            if t_elem is None:
                t_elem = etree.SubElement(new_run, f'{{{_NS_A}}}t')
            t_elem.text = seg_text
        p.insert(insert_pos, new_run)
        insert_pos += 1


def _clone_cell_paragraph(src_p, src_run, sz):
    """양식 셀의 pPr(정렬·줄간격)·run rPr(Pretendard·색) 보존한 빈 문단 생성.
    sz만 지정값으로 교체. 텍스트 빈 run 1개 포함."""
    p = etree.Element(f'{{{_NS_A}}}p')
    src_pPr = src_p.find(f'{{{_NS_A}}}pPr')
    if src_pPr is not None:
        p.append(copy.deepcopy(src_pPr))        # 정렬 ctr·lnSpc 100% 그대로
    run = copy.deepcopy(src_run)                # Pretendard·색 그대로
    rPr = run.find(f'{{{_NS_A}}}rPr')
    if rPr is not None:
        rPr.set('sz', sz)                       # ← 폰트 크기만 교체
    t = run.find(f'{{{_NS_A}}}t')
    if t is None:
        t = etree.SubElement(run, f'{{{_NS_A}}}t')
    t.text = ''
    p.append(run)
    return p


def _make_date_paragraph(date_num, color=DATE_COLOR_GREEN):
    """날짜 전용 paragraph 생성: 왼쪽 정렬, 9pt, 볼드, Pretendard.
    color: 평일·공휴일·선거일 모두 #318E72(원본 정답본 초록)으로 통일."""
    date_p = etree.Element(f'{{{_NS_A}}}p')

    pPr = etree.SubElement(date_p, f'{{{_NS_A}}}pPr')
    pPr.set('algn', 'l')

    rPr = etree.Element(f'{{{_NS_A}}}rPr')
    rPr.set('lang', 'ko-KR')
    rPr.set('sz', '900')
    rPr.set('b', '1')
    rPr.set('dirty', '0')

    solidFill = etree.SubElement(rPr, f'{{{_NS_A}}}solidFill')
    srgbClr   = etree.SubElement(solidFill, f'{{{_NS_A}}}srgbClr')
    srgbClr.set('val', color)

    latin = etree.SubElement(rPr, f'{{{_NS_A}}}latin')
    latin.set('typeface', 'Pretendard')
    ea = etree.SubElement(rPr, f'{{{_NS_A}}}ea')
    ea.set('typeface', 'Pretendard')

    run = etree.SubElement(date_p, f'{{{_NS_A}}}r')
    run.append(rPr)
    t = etree.SubElement(run, f'{{{_NS_A}}}t')
    t.text = date_num

    return date_p


def _make_reason_paragraph(reason_text):
    """공휴일 사유 단락 생성: 왼쪽 정렬, #318E72, 8pt, 볼드 없음, Pretendard."""
    p = etree.Element(f'{{{_NS_A}}}p')

    pPr = etree.SubElement(p, f'{{{_NS_A}}}pPr')
    pPr.set('algn', 'l')

    rPr = etree.Element(f'{{{_NS_A}}}rPr')
    rPr.set('lang', 'ko-KR')
    rPr.set('sz', '800')
    rPr.set('dirty', '0')

    solidFill = etree.SubElement(rPr, f'{{{_NS_A}}}solidFill')
    srgbClr   = etree.SubElement(solidFill, f'{{{_NS_A}}}srgbClr')
    srgbClr.set('val', DATE_COLOR_GREEN)

    latin = etree.SubElement(rPr, f'{{{_NS_A}}}latin')
    latin.set('typeface', 'Pretendard')
    ea = etree.SubElement(rPr, f'{{{_NS_A}}}ea')
    ea.set('typeface', 'Pretendard')

    run = etree.SubElement(p, f'{{{_NS_A}}}r')
    run.append(rPr)
    t = etree.SubElement(run, f'{{{_NS_A}}}t')
    t.text = reason_text

    return p


def set_lunch_cell(cell, lines, date_num=None):
    """본식 셀: 원본 정답 서식(항목당 1문단·8.5pt·br없음·가운데정렬)으로 재생성.
    lines: [밥, 국, 반찬1, ..., 영양구성]. date_num 있으면 날짜 문단(9pt) 맨 앞 유지."""
    tc     = cell._tc
    txBody = tc.find(f'{{{_NS_A}}}txBody')
    if txBody is None:
        return
    src_p = txBody.find(f'{{{_NS_A}}}p')
    if src_p is None:
        return
    src_run = src_p.find(f'{{{_NS_A}}}r')
    if src_run is None:
        return
    # 기존 문단 전부 제거 (양식 run-br 틀 폐기 후 재생성)
    for p in txBody.findall(f'{{{_NS_A}}}p'):
        txBody.remove(p)
    # 날짜 문단 (우리 구조 유지, 9pt — 변경 없음)
    if date_num:
        txBody.append(_make_date_paragraph(date_num))
    # 메뉴 항목별 문단 (8.5pt, 인라인 알레르기 빨강)
    for text in lines:
        if not text:
            continue
        p = _clone_cell_paragraph(src_p, src_run, '850')
        _apply_inline_allergy(p.find(f'{{{_NS_A}}}r'), text, p)
        txBody.append(p)


def _set_date_only_cell(cell, date_num, reason=''):
    """빈 셀(공휴일 clear 후)에 날짜(+사유) paragraph 삽입.
    날짜·사유 모두 #318E72(원본 초록)로 통일 — 공휴일/평일 색 구분 없음.
    공휴일 칸은 내용 적어 기본(ctr)이면 날짜가 중앙으로 처짐 →
    tcPr anchor=t로 위쪽 정렬, 평일 날짜와 높이 맞춤.
    현충일 등 다른 공휴일도 동일하게 위쪽 정렬 — 의도된 일관 동작."""
    tc     = cell._tc
    txBody = tc.find(f'{{{_NS_A}}}txBody')
    if txBody is None or not date_num:
        return
    date_p = _make_date_paragraph(date_num, color=DATE_COLOR_GREEN)
    existing_p = txBody.find(f'{{{_NS_A}}}p')
    if existing_p is not None:
        p_idx = list(txBody).index(existing_p)
        txBody.insert(p_idx, date_p)
        if reason:
            txBody.insert(p_idx + 1, _make_reason_paragraph(reason))
    else:
        txBody.append(date_p)
        if reason:
            txBody.append(_make_reason_paragraph(reason))
    # 표 셀 수직정렬은 tcPr.anchor (bodyPr 아님)
    tcPr = tc.find(f'{{{_NS_A}}}tcPr')
    if tcPr is None:
        tcPr = etree.SubElement(tc, f'{{{_NS_A}}}tcPr')
    tcPr.set('anchor', 't')


def set_snack_cell(cell, menu_line, kcal_line):
    """간식 셀: 원본 정답 서식(2문단 메뉴/Kcal·6.5pt·br없음·가운데정렬)으로 재생성."""
    tc     = cell._tc
    txBody = tc.find(f'{{{_NS_A}}}txBody')
    if txBody is None:
        return
    src_p = txBody.find(f'{{{_NS_A}}}p')
    if src_p is None:
        return
    src_run = src_p.find(f'{{{_NS_A}}}r')
    if src_run is None:
        return
    for p in txBody.findall(f'{{{_NS_A}}}p'):
        txBody.remove(p)
    # 문단0: 메뉴 (6.5pt, 인라인 알레르기 빨강)
    p_menu = _clone_cell_paragraph(src_p, src_run, '650')
    _apply_inline_allergy(p_menu.find(f'{{{_NS_A}}}r'), menu_line, p_menu)
    txBody.append(p_menu)
    # 문단1: Kcal (6.5pt)
    if kcal_line:
        p_kcal = _clone_cell_paragraph(src_p, src_run, '650')
        _set_run_text(p_kcal.find(f'{{{_NS_A}}}r'), kcal_line)
        txBody.append(p_kcal)


def set_label_cell(cell, label):
    _, runs = _get_p_and_runs(cell)
    if runs:
        _set_run_text(runs[0], label)


def clear_cell(cell):
    """셀 내 모든 run·br 제거 + 여분 단락 제거 → 완전 빈 셀."""
    tc     = cell._tc
    txBody = tc.find(f'{{{_NS_A}}}txBody')
    if txBody is None:
        return
    all_ps = txBody.findall(f'{{{_NS_A}}}p')
    for extra_p in all_ps[1:]:
        txBody.remove(extra_p)
    p = txBody.find(f'{{{_NS_A}}}p')
    if p is None:
        return
    for child in list(p):
        tag = child.tag
        if tag == f'{{{_NS_A}}}r' or tag == f'{{{_NS_A}}}br':
            p.remove(child)


# ════════════════════════════════════════════════════════════════════
# 6. 메뉴 → 셀 라인 빌더
# ════════════════════════════════════════════════════════════════════
_LUNCH_KEYS = ['bap', 'guk', 'banchan1', 'banchan2', 'banchan3', 'banchan4', 'banchan5']


def build_lunch_lines(lunch, branch, do_strip=False, add_fruit=False, fruit_text='제철과일'):
    lines = []
    for key in _LUNCH_KEYS:
        resolved = resolve_for_branch(lunch.get(key, ''), branch)
        if resolved and resolved.strip():
            if do_strip:
                resolved = strip_allergy(resolved)
            lines.append(resolved.strip())

    if add_fruit and lunch.get('has_fruit'):
        lines.append(fruit_text)

    nutri = lunch.get('nutrition', '')
    if nutri and nutri.strip():
        lines.append(nutri.strip())

    return lines


def build_snack(snack, branch, do_strip=False):
    menu = resolve_for_branch(snack.get('menu', ''), branch)
    kcal = snack.get('nutrition', '')
    if do_strip:
        menu = strip_allergy(menu)
    return menu.strip(), kcal.strip()


# ════════════════════════════════════════════════════════════════════
# 7. 슬라이드 메타데이터 교체
# ════════════════════════════════════════════════════════════════════
def _iter_all_shapes(shapes):
    for s in shapes:
        yield s
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                yield from _iter_all_shapes(s.shapes)
            except Exception:
                pass


def _replace_in_tf(tf, find, replace):
    for para in tf.paragraphs:
        runs = para.runs
        if not runs:
            continue
        for run in runs:
            if find in run.text:
                run.text = run.text.replace(find, replace)
                return True
        combined = ''.join(r.text for r in runs)
        if find in combined:
            runs[0].text = combined.replace(find, replace)
            for r in runs[1:]:
                r.text = ''
            return True
    return False


def replace_slide_metadata(slide, display_name, email):
    for shape in _iter_all_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        _replace_in_tf(tf, 'OOO 지점', f'{display_name} 지점')
        _replace_in_tf(tf, 'OOO 영양사', '영양사')
        _replace_in_tf(tf, 'OOOOOOO@kizmeal.com', email)


_DATE_NUM_RE = re.compile(r'^(0[1-9]|[12][0-9]|3[01])$')


def replace_date_group(slide, date_map=None):
    """
    lxml으로 spTree를 직접 순회해 전체 텍스트가 '01'~'31'인
    p:sp를 완전 삭제 (그룹 중첩 깊이 무관).
    """
    spTree = slide.shapes._spTree
    _SP_TAG = f'{{{_NS_PML}}}sp'
    _T_TAG  = f'{{{_NS_A}}}t'

    # 먼저 삭제 대상 수집 후 제거 (순회 중 remove 방지)
    to_remove = []
    for sp in spTree.findall(f'.//{_SP_TAG}'):
        txt = ''.join(t.text or '' for t in sp.findall(f'.//{_T_TAG}')).strip()
        if _DATE_NUM_RE.match(txt):
            to_remove.append(sp)

    for sp in to_remove:
        parent = sp.getparent()
        if parent is not None:
            parent.remove(sp)


# 양식 라벨 제거 시 절대 보존할 도형 이름 (코드가 이름으로 직접 참조)
_LABEL_PRESERVE_NAMES = frozenset({'ORIGIN_BOX', 'MATERIAL_BOX', 'ALLERGY_BOX'})


def remove_holiday_labels(slide, reasons):
    """양식에 떠있는 공휴일/선거일 라벨 TextBox를 안전 제거.
    우리가 셀에 날짜+사유를 직접 넣으므로 양식의 떠있는 라벨은 중복.

    3중 가드 (이름 의존 금지 · 동적):
      1) 단독 TEXT_BOX (sp) — pic/group/table 제외
      2) 최상위 도형 (slide.shapes 직접 순회 = 그룹 내부 미진입, 헤더 장식 보호)
      3) 이름이 보존 화이트리스트에 없음 (ORIGIN/MATERIAL/ALLERGY 보호)
      4) 좌표(좌상단)가 MENU_TABLE bbox 내부 — 표 밖 범례/제목 보호
      5) 텍스트 매칭: reason in 도형텍스트
         (양식라벨 '지방선거일' 이 우리 reason '선거일' 을 포함 — 방향 주의)
    수집 후 일괄 remove (순회 중 변형 방지 — replace_date_group 과 동일).
    """
    if not reasons:
        return

    # MENU_TABLE 도형 bbox 확보 (이름 우선, 없으면 첫 표)
    tbl_shape = None
    for sh in slide.shapes:
        try:
            if sh.name == 'MENU_TABLE' and sh.has_table:
                tbl_shape = sh
                break
        except Exception:
            continue
    if tbl_shape is None:
        for sh in slide.shapes:
            try:
                if sh.has_table:
                    tbl_shape = sh
                    break
            except Exception:
                continue
    if tbl_shape is None:
        return

    t_left   = tbl_shape.left or 0
    t_top    = tbl_shape.top or 0
    t_right  = t_left + (tbl_shape.width or 0)
    t_bottom = t_top + (tbl_shape.height or 0)

    to_remove = []
    for sh in slide.shapes:                       # 최상위만 (그룹 내부 미진입)
        try:
            if sh.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
                continue
        except Exception:
            continue
        if sh.name in _LABEL_PRESERVE_NAMES:
            continue
        try:
            if not sh.has_text_frame:
                continue
        except Exception:
            continue
        l = sh.left or 0
        t = sh.top or 0
        if not (t_left <= l <= t_right and t_top <= t <= t_bottom):
            continue
        txt = sh.text_frame.text
        if not any(r and r in txt for r in reasons):
            continue
        to_remove.append(sh)

    for sh in to_remove:
        print(f"  🗑️  양식 라벨 제거: name='{sh.name}' "
              f"text={sh.text_frame.text.strip()!r}")
        sp = sh._element
        parent = sp.getparent()
        if parent is not None:
            parent.remove(sp)


def _holiday_cell_rects(tbl_shape, sections, week_days):
    """공휴일 날짜(is_holiday) 칸의 영역 rect 목록 반환.
    각 rect = (left, top, right, bottom) EMU.
    세로는 해당 주의 전체 섹션 행을 합친 범위(점심~오후)로 잡아
    날짜·메뉴 어디를 가리든 겹침 판정되게 한다.
    행 인덱싱은 _render_slide 와 동일: row = (wk-1)*len(sections) + sec순번."""
    tbl    = tbl_shape.table
    cols   = [c.width  for c in tbl.columns]
    rows   = [r.height for r in tbl.rows]
    t_left = tbl_shape.left or 0
    t_top  = tbl_shape.top or 0

    # 열·행 누적 경계 좌표
    col_x = [t_left]
    for w in cols:
        col_x.append(col_x[-1] + w)
    row_y = [t_top]
    for h in rows:
        row_y.append(row_y[-1] + h)

    n_sec = len(sections)
    rects = []
    for wk in range(1, 6):
        days = week_days.get(wk, [None] * 5)
        for col in range(5):
            day = days[col]
            if not (day and day.get('is_holiday')):
                continue
            cell_col = col + 1          # col0 = 라벨열
            if cell_col + 1 >= len(col_x):
                continue
            r0 = (wk - 1) * n_sec
            r1 = min(r0 + n_sec, len(rows))
            if r0 >= len(rows):
                continue
            rects.append((col_x[cell_col], row_y[r0],
                          col_x[cell_col + 1], row_y[r1]))
    return rects


def remove_holiday_cover_pics(slide, sections, week_days):
    """공휴일 셀 위에 떠서 날짜를 가리는 양식 장식 PICTURE(예: 선거 도장 그림4)를
    동적 제거. B-2 라벨 제거의 그림(좌표) 버전.

    제거 조건 (전부 AND):
      1) 최상위 PICTURE (slide.shapes 직접 순회 = 그룹 내부 미진입)
      2) MENU_TABLE 보다 위 z-order (spTree 순서상 뒤 = 위)
      3) bbox 가 공휴일 셀 영역과 겹침 (그림은 텍스트 없음 → 좌표 겹침으로 판정)
    보존: 로고(그림27)·하단 클립보드(그림3)는 셀과 안 겹쳐 자동 보존.
          알레르기/원산지 박스는 PICTURE 아님 → 미대상.
    수집 후 일괄 remove.
    """
    # 표 도형 확보
    tbl_shape = None
    for sh in slide.shapes:
        try:
            if sh.name == 'MENU_TABLE' and sh.has_table:
                tbl_shape = sh
                break
        except Exception:
            continue
    if tbl_shape is None:
        for sh in slide.shapes:
            try:
                if sh.has_table:
                    tbl_shape = sh
                    break
            except Exception:
                continue
    if tbl_shape is None:
        return

    rects = _holiday_cell_rects(tbl_shape, sections, week_days)
    if not rects:
        return

    spTree   = slide.shapes._spTree
    children = list(spTree)
    try:
        tbl_idx = children.index(tbl_shape._element)
    except ValueError:
        tbl_idx = -1

    to_remove = []
    for sh in slide.shapes:                       # 최상위만
        try:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            continue
        try:
            z = children.index(sh._element)
        except ValueError:
            continue
        if z <= tbl_idx:                          # 표보다 아래면 가리지 않음
            continue
        pl = sh.left or 0
        pt = sh.top or 0
        pr = pl + (sh.width or 0)
        pb = pt + (sh.height or 0)
        hit = any(not (pr <= rl or pl >= rr or pb <= rt or pt >= rb)
                  for (rl, rt, rr, rb) in rects)
        if hit:
            to_remove.append(sh)

    for sh in to_remove:
        print(f"  🗑️  공휴일칸 가림 그림 제거: name='{sh.name}' "
              f"pos=({sh.left:,},{sh.top:,}) size=({sh.width:,}x{sh.height:,})")
        pic = sh._element
        parent = pic.getparent()
        if parent is not None:
            parent.remove(pic)


# ════════════════════════════════════════════════════════════════════
# 8. 표 헬퍼
# ════════════════════════════════════════════════════════════════════
def _get_table(slide):
    # 1순위: MENU_TABLE 이름표로 찾기
    for shape in slide.shapes:
        try:
            if shape.name == 'MENU_TABLE' and shape.has_table:
                return shape.table
        except Exception:
            continue
    # 2순위: 첫 번째 has_table 폴백 (이름표 없는 기존 템플릿 대응)
    for shape in slide.shapes:
        try:
            if shape.has_table:
                return shape.table
        except Exception:
            continue
    return None


def _days_by_week(menu_data):
    out   = {}
    weeks = menu_data.get('weeks', {})
    for wk in range(1, 6):
        wdata = weeks.get(str(wk))
        if not wdata:
            out[wk] = [None] * 5
            continue
        days   = wdata.get('days', [])
        padded = []
        for d in days[:5]:
            padded.append(None if (d is None or d.get('is_skipped')) else d)
        padded += [None] * (5 - len(padded))
        out[wk] = padded
    return out


# ════════════════════════════════════════════════════════════════════
# 8-1. 이름표(cNvPr name) 기반 도형 탐색
# ════════════════════════════════════════════════════════════════════
def find_shape_by_name(slide_el, target_name):
    """
    이름표(cNvPr name)로 도형 찾기.
    sp(텍스트박스), graphicFrame(표), grpSp(그룹) 모두 탐색.
    ALLERGY_BOX는 그룹이라 grpSp까지 봐야 함.
    찾으면 element 반환, 없으면 None.
    """
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    for tag in ('sp', 'graphicFrame', 'grpSp'):
        for el in slide_el.iter(f'{{{ns_p}}}{tag}'):
            cNvPr = el.find(f'.//{{{ns_p}}}cNvPr')
            if cNvPr is not None and cNvPr.get('name') == target_name:
                return el
    return None


def find_shape_smart(slide_el, target_name, fallback_keyword=None):
    """
    이름표 우선 탐색, 없으면 텍스트 키워드로 폴백.
    기존 템플릿(이름표 없음)과 새 템플릿(이름표 있음) 모두 지원.
    """
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # 1순위: 이름표로 찾기
    el = find_shape_by_name(slide_el, target_name)
    if el is not None:
        return el, 'name'

    # 2순위: 텍스트 키워드로 폴백
    if fallback_keyword:
        for sp in slide_el.iter(f'{{{ns_p}}}sp'):
            texts = ''.join(t.text or '' for t in sp.findall(f'.//{{{ns_a}}}t'))
            if fallback_keyword in texts:
                return sp, 'keyword'

    return None, None


# ════════════════════════════════════════════════════════════════════
# 9. 슬라이드 렌더링
# ════════════════════════════════════════════════════════════════════
def _render_slide(table, plan_entry, week_days, cfg):
    _, sections, opts = plan_entry
    if opts.get('skip') or not sections:
        return

    do_strip    = opts.get('strip', False)
    care_label  = opts.get('care_label')
    snack_label = cfg.get('snack_label', '오전')
    add_fruit   = cfg.get('add_fruit', False)
    fruit_text  = cfg.get('fruit_text', '제철과일')
    fixed_am    = cfg.get('fixed_am')
    use_pm_as_am= cfg.get('use_pm_as_am', False)
    branch_name = cfg['name']

    _init_table_cells(table, sections)

    n_rows  = len(table.rows)
    row_idx = 0

    for wk in range(1, 6):
        days = week_days.get(wk, [None] * 5)
        for sec in sections:
            if row_idx >= n_rows:
                break
            row   = table.rows[row_idx]
            cells = row.cells
            row_idx += 1

            if sec == 'am' and snack_label == '간식':
                set_label_cell(cells[0], '간식')
            elif sec == 'care' and care_label:
                set_label_cell(cells[0], care_label)

            for col in range(5):
                if col + 1 >= len(cells):
                    break
                cell = cells[col + 1]
                day  = days[col]

                if day is None:
                    clear_cell(cell)
                    continue

                if sec == 'lunch':
                    date_str = day.get('date', '')
                    # ISO '2026-06-03' 또는 2자리 '03' 둘 다 처리, 앞의 0 제거
                    if len(date_str) == 10 and date_str[4] == '-':
                        day_num = date_str.split('-')[2]
                    elif date_str.isdigit():
                        day_num = date_str
                    else:
                        day_num = ''
                    day_num = str(int(day_num)) if day_num else ''
                    if day.get('is_holiday') and branch_name not in _HOLIDAY_OPERATING:
                        clear_cell(cell)
                        _set_date_only_cell(cell, day_num, reason=day.get('holiday_reason', ''))
                        continue
                    lines = build_lunch_lines(
                        day.get('lunch', {}), branch_name,
                        do_strip=do_strip,
                        add_fruit=add_fruit,
                        fruit_text=fruit_text,
                    )
                    if not lines:
                        # 운영원 공휴일이지만 데이터 없는 경우 — 날짜만 표시
                        clear_cell(cell)
                        _set_date_only_cell(cell, day_num)
                        continue
                    set_lunch_cell(cell, lines, date_num=day_num)

                elif sec == 'am':
                    if day.get('is_holiday') and branch_name not in _HOLIDAY_OPERATING:
                        clear_cell(cell)
                        continue
                    if fixed_am:
                        set_snack_cell(cell, fixed_am['menu'], fixed_am['nutrition'])
                        continue
                    snack_key = 'afternoon_snack' if use_pm_as_am else 'morning_snack'
                    menu, kcal = build_snack(
                        day.get(snack_key, {}), branch_name, do_strip=do_strip)
                    if not menu:
                        clear_cell(cell)
                        continue
                    set_snack_cell(cell, menu, kcal)

                elif sec == 'pm':
                    if day.get('is_holiday') and branch_name not in _HOLIDAY_OPERATING:
                        clear_cell(cell)
                        continue
                    menu, kcal = build_snack(
                        day.get('afternoon_snack', {}), branch_name, do_strip=do_strip)
                    if not menu:
                        clear_cell(cell)
                        continue
                    set_snack_cell(cell, menu, kcal)

                elif sec == 'care':
                    if day.get('is_holiday') and branch_name not in _HOLIDAY_OPERATING:
                        clear_cell(cell)
                        continue
                    menu, kcal = build_snack(
                        day.get('care_snack', {}), branch_name, do_strip=do_strip)
                    set_snack_cell(cell, menu, kcal)


# ════════════════════════════════════════════════════════════════════
# 9-2. 텍스트박스 교체
# ════════════════════════════════════════════════════════════════════
_MATERIAL_FIXED_LINE = '키즈밀은 엄선된 친환경 및 국내산 식자재를 선별하여 사용하고, 화학첨가물이 들어있는 가공식품은 최대한 배제하여 올바른 식습관 확립 및 균형 잡힌 영양을 제공합니다.'

_SP_TAG_PML = f'{{{_NS_PML}}}sp'
_P_TAG      = f'{{{_NS_A}}}p'
_T_TAG_A    = f'{{{_NS_A}}}t'
_BODY_PR    = f'{{{_NS_A}}}bodyPr'


def _make_text_paragraph(text, sz=900, bold=False, algn=None):
    """단순 텍스트 paragraph XML 생성."""
    p = etree.Element(_P_TAG)
    if algn:
        pPr = etree.SubElement(p, f'{{{_NS_A}}}pPr')
        pPr.set('algn', algn)
    rPr = etree.Element(f'{{{_NS_A}}}rPr')
    rPr.set('lang', 'ko-KR')
    rPr.set('sz', str(sz))
    if bold:
        rPr.set('b', '1')
    rPr.set('dirty', '0')
    run = etree.SubElement(p, f'{{{_NS_A}}}r')
    run.append(rPr)
    t = etree.SubElement(run, f'{{{_NS_A}}}t')
    t.text = text or ''
    return p


# p:sp 안의 txBody는 p:(PML) 또는 a:(DrawingML) 네임스페이스 둘 다 가능
_TXBODY_TAGS = [f'{{{_NS_PML}}}txBody', f'{{{_NS_A}}}txBody']


def _find_txbody(sp):
    """p:txBody → a:txBody 순으로 탐색."""
    for tag in _TXBODY_TAGS:
        txBody = sp.find(f'.//{tag}')
        if txBody is not None:
            return txBody
    return None


def _replace_textbox_content(slide, search_text, paragraphs, target_name=None):
    spTree = slide.shapes._spTree
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    target_sp = None
    found_by = None

    # 1순위: 이름표(target_name)로 찾기
    if target_name:
        for sp in spTree.findall(f'.//{_SP_TAG_PML}'):
            cNvPr = sp.find(f'.//{{{ns_p}}}cNvPr')
            if cNvPr is not None and cNvPr.get('name') == target_name:
                target_sp = sp
                found_by = 'name'
                break

    # 2순위: 텍스트 키워드로 폴백 (이름표 없는 기존 템플릿 대응)
    if target_sp is None:
        for sp in spTree.findall(f'.//{_SP_TAG_PML}'):
            full_text = ''.join(t.text or '' for t in sp.findall(f'.//{_T_TAG_A}'))
            if search_text in full_text:
                target_sp = sp
                found_by = 'keyword'
                break

    if target_sp is None:
        return

    if target_name:
        print(f"  {target_name} 탐색: {found_by} 방식")

    txBody = _find_txbody(target_sp)
    if txBody is None:
        return
    for p in txBody.findall(_P_TAG):
        txBody.remove(p)
    for para_cfg in paragraphs:
        new_p = _make_text_paragraph(
            para_cfg.get('text', ''),
            sz=para_cfg.get('sz', 900),
            bold=para_cfg.get('bold', False),
            algn=para_cfg.get('algn'),
        )
        txBody.append(new_p)
    return


def replace_origin_text(slide, origin_text):
    """원산지 텍스트박스 교체."""
    if not origin_text:
        return
    _replace_textbox_content(slide, '원산지 표기', [
        {'text': '원산지 표기', 'sz': 1000, 'bold': True, 'algn': 'l'},
        {'text': origin_text.get('body', ''),       'sz': 800},
        {'text': origin_text.get('disclaimer', ''), 'sz': 800},
    ], target_name='ORIGIN_BOX')


def replace_material_text(slide, material_text):
    """원재료 텍스트박스 교체."""
    if not material_text:
        return
    _replace_textbox_content(slide, '원재료 표시안내', [
        {'text': '* 원재료 표시안내 *', 'sz': 800, 'bold': False, 'algn': 'l'},
        {'text': material_text,         'sz': 800},
        {'text': _MATERIAL_FIXED_LINE,  'sz': 800},
    ], target_name='MATERIAL_BOX')


# ════════════════════════════════════════════════════════════════════
# 9-3. 15행 행높이 정규화 + 텍스트 크기 보정 + 박스 겹침 감지
# ════════════════════════════════════════════════════════════════════
def _fix_allergy_only(prs):
    """ORIGIN_BOX/MATERIAL_BOX/ALLERGY_BOX **위치는 손대지 않는다.**

    ⚠️ 2026-08-11 변경: 예전엔 이 함수가 세 박스 위치를 하드코딩된 오프셋으로
    강제 이동시켰다(ALLERGY_BOX는 table_bottom/슬라이드 하단 기준 재계산,
    ORIGIN/MATERIAL은 +175,260 EMU 하향 또는 8,395,899로 고정). 로컬 6월
    템플릿(당시 유일하게 쓰이던 템플릿)의 박스 위치가 어긋나 있던 걸 코드로
    보정한 것이었는데, 이름표 자동부여(발견①) 이후 매달 진짜 디자이너 원본을
    쓰게 되면서 디자이너가 그 달 빈칸 위치에 박스를 이미 배치해서 준다는 게
    실측으로 확인됐다(HANDOFF "식단표 자동화" 조사 발견②). 실측: 8월 양식
    ORIGIN top=8,428,017인데 이 함수는 여기서 175,260을 더 밀어 8,603,277로
    만들고 있었음 — 실제 정답 배포본(8,450,877)과도 다른 값. 박스 위치 이동
    로직 전부 제거, 디자이너가 준 위치를 그대로 신뢰한다.

    ⚠️ 2026-08-12 변경: 15행 슬라이드의 행 높이 강제 재설정(_H_LUNCH_15/
    _H_SNACK_15)과 graphicFrame ext.cy 재계산도 제거했다. 영양사는 행
    높이를 지정하지 않는다 — 기준 엑셀 식단표에서 중식·오전/오후 간식 칸을
    드래그 복사해 디자이너 PPTX 표에 붙여넣는 방식이며, 그 과정에서 엑셀
    쪽 행 높이가 딸려온다(영양사 실제 편집 과정 확인, 2026-08-12). 즉
    배포본의 행 높이는 사양이 아니라 복사-붙여넣기의 부산물이고, 그 달
    기준 엑셀에 따라 매달 달라진다(실측: 6월 실물 1,180,800/306,000,
    7월은 디자이너 값 그대로 제각각, 8월은 1,152,000 균일 — 세 달 연속
    서로 다름, 규칙 없음). 코드가 재현해야 할 대상이 아니다.
    지웠던 상수(_H_LUNCH_15=1,180,800 / _H_SNACK_15=306,000)는
    `pptx-server/templates/(강동ECC).pptx`에서 역산한 값이다. 이 파일은
    6월 실제 배포본과 SHA-256이 완전히 동일한 진짜 6월 정답본이며(오염
    fixture 아님, 확인 완료), 문제는 파일 자체가 아니라 **"6월 한 달의
    실측값을 상수로 굳혀 7·8월 등 다른 달 생성에도 그대로 적용한 것"**이다
    — 달마다 다른 값을 영구 상수로 하드코딩한 것이 잘못이었다. 이제
    디자이너 원본의 a:tr/@h를 그대로 통과시키고 ext.cy도 손대지 않는다
    (실물도 PowerPoint가 ext.cy를 갱신하지 않은 상태로 배포되므로, 원본을
    그대로 두는 편이 오히려 실물에 더 가깝다).

    남은 역할(10행 경로에서만 수행 — 15행은 위에서 바로 continue): (1)
    박스 텍스트 폰트 크기 보정(위치는 안 건드림), (2) 로고·클립보드
    아이콘 리사이즈(크기만, 위치 이동은 제거 — 클립보드는 원래 ORIGIN_BOX
    위치 기준으로 나란히 정렬하던 것이라 그 기준점이 없어져 위치 보정도
    함께 뺌. 로고는 원래도 자기 자신 기준 우측/하단 정렬이라 영향 없음,
    그대로 둠).
    """
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    for slide in prs.slides:
        slide_el = slide._element

        # MENU_TABLE 이름표 우선으로 테이블 찾기 (없으면 첫 graphicFrame 폴백)
        table_el, _ = find_shape_smart(slide_el, 'MENU_TABLE', fallback_keyword=None)
        if table_el is None:
            table_el = slide_el.find(f'.//{{{ns_p}}}graphicFrame')
        if table_el is None:
            continue

        # graphicFrame의 xfrm은 PML 네임스페이스 (p:xfrm), off/ext는 DrawingML
        p_xfrm = table_el.find(f'{{{ns_p}}}xfrm')
        tbl = table_el.find(f'.//{{{ns_a}}}tbl')
        if p_xfrm is None or tbl is None:
            continue
        rows = tbl.findall(f'{{{ns_a}}}tr')
        row_count = len(rows)

        if row_count >= 15:
            # 2026-08-12: 행 높이 강제 재설정(_H_LUNCH_15/_H_SNACK_15)과
            # graphicFrame ext.cy 동기화를 제거. 디자이너 원본의 a:tr/@h,
            # ext.cy를 그대로 통과시킨다 — 근거는 위 함수 docstring 참고.
            continue

        # 10행: ALLERGY_BOX 폰트만 보정(위치 이동 제거)
        sp, found_by = find_shape_smart(slide_el, 'ALLERGY_BOX', fallback_keyword='알레르기 표시')
        if sp is not None:
            for rPr in sp.findall(f'.//{{{ns_a}}}rPr'):
                szv = rPr.get('sz')
                if szv is None or int(szv) >= 700:
                    rPr.set('sz', '500')
            print(f"  ✅ ALLERGY_BOX 폰트5pt ({found_by}, 위치 불변)")

        # 원산지/원재료 본문 폰트 크기만 보정 (위치는 손대지 않음)
        for box_name in ('ORIGIN_BOX', 'MATERIAL_BOX'):
            bsp, _ = find_shape_smart(slide_el, box_name, fallback_keyword=None)
            if bsp is None:
                continue
            for rPr in bsp.findall(f'.//{{{ns_a}}}rPr'):
                szv = rPr.get('sz')
                if szv is None:
                    rPr.set('sz', '550')
                else:
                    s = int(szv)
                    if s >= 1000:
                        rPr.set('sz', '850')   # 제목 8.5pt 유지
                    elif s >= 600:
                        rPr.set('sz', '550')   # 본문 5.5pt
            print(f"  ✅ {box_name} 폰트5.5pt (위치 불변)")

        # ② 키즈밀 로고 축소 + 하단 이동 (pic name='그래픽 118')
        SLIDE_H = 10680700
        LOGO_BOTTOM_MARGIN = 300000  # 슬라이드 하단에서 띄울 여백
        for pic in slide_el.findall(f'.//{{{ns_p}}}pic'):
            c = pic.find(f'.//{{{ns_p}}}cNvPr')
            if c is None or c.get('name') != '그래픽 118':
                continue
            xfrm = pic.find(f'.//{{{ns_a}}}xfrm')
            if xfrm is None:
                continue
            off = xfrm.find(f'{{{ns_a}}}off')
            ext = xfrm.find(f'{{{ns_a}}}ext')
            if off is None or ext is None:
                continue
            ox = int(off.get('x')); oy = int(off.get('y'))
            ocx = int(ext.get('cx')); ocy = int(ext.get('cy'))
            # 80% 축소 (가로세로 비율 유지, 우측 정렬 유지)
            ncx = int(ocx * 0.8); ncy = int(ocy * 0.8)
            right_edge = ox + ocx
            nx = right_edge - ncx        # 우측 끝 맞춤
            ny = SLIDE_H - ncy - LOGO_BOTTOM_MARGIN  # 하단에 붙임
            ext.set('cx', str(ncx)); ext.set('cy', str(ncy))
            off.set('x', str(nx)); off.set('y', str(ny))
            print(f"  ✅ 로고 축소80%+하단이동: y={ny:,} cy={ncy:,}")
            break

        # 식단표 클립보드 아이콘 리사이즈만(위치는 그대로) — "원산지 제목과 나란히"
        # 정렬하던 위치 보정은 ORIGIN_BOX를 더 이상 옮기지 않게 되면서 기준점을
        # 잃어 함께 제거함(2026-08-11, 박스 이동 로직 제거).
        for pic in slide_el.findall(f'.//{{{ns_p}}}pic'):
            c = pic.find(f'.//{{{ns_p}}}cNvPr')
            if c is None or c.get('name') not in {'그림 3', '그림 224', '그림 193'}:
                continue
            xfrm = pic.find(f'.//{{{ns_a}}}xfrm')
            if xfrm is None:
                continue
            off = xfrm.find(f'{{{ns_a}}}off')
            ext = xfrm.find(f'{{{ns_a}}}ext')
            if off is None or ext is None:
                continue
            ocx = int(ext.get('cx')); ocy = int(ext.get('cy'))
            # 80% 축소
            ncx = int(ocx * 0.8); ncy = int(ocy * 0.8)
            ext.set('cx', str(ncx)); ext.set('cy', str(ncy))
            print(f"  ✅ 클립보드 아이콘 80% 축소 (위치 불변)")
            break


def _check_box_overlap(prs, branch_label=''):
    """원산지/원재료/알레르기 세 박스가 겹치는지 검사, 겹치면 경고 출력."""
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    def get_rect(slide_el, name):
        sp, _ = find_shape_smart(slide_el, name, fallback_keyword=None)
        if sp is None:
            return None
        xfrm = sp.find(f'.//{{{ns_a}}}xfrm')
        if xfrm is None:
            return None
        off = xfrm.find(f'{{{ns_a}}}off')
        ext = xfrm.find(f'{{{ns_a}}}ext')
        if off is None or ext is None:
            return None
        x  = int(off.get('x'));  y  = int(off.get('y'))
        cx = int(ext.get('cx')); cy = int(ext.get('cy'))
        return (x, y, x + cx, y + cy)

    def overlap(a, b):
        if a is None or b is None:
            return False
        ax1, ay1, ax2, ay2 = a
        bx1, by1, bx2, by2 = b
        return not (ax2 <= bx1 or bx2 <= ax1 or ay2 <= by1 or by2 <= ay1)

    names = ['ORIGIN_BOX', 'MATERIAL_BOX', 'ALLERGY_BOX']
    for si, slide in enumerate(prs.slides):
        rects = {n: get_rect(slide._element, n) for n in names}
        for i in range(len(names)):
            for j in range(i + 1, len(names)):
                if overlap(rects[names[i]], rects[names[j]]):
                    print(f"  ⚠️ 겹침 감지! {branch_label} 슬라이드{si + 1}: {names[i]} ↔ {names[j]}")


# ════════════════════════════════════════════════════════════════════
# 10. 메인 생성 함수
# ════════════════════════════════════════════════════════════════════
def generate(cfg, menu_data, template_path, out_path, date_map=None,
             origin_text=None, material_text=None):
    type_code = cfg['type']
    plan = TYPE_PLANS.get(type_code)
    if plan is None:
        raise ValueError(f'알 수 없는 타입: {type_code}')

    slide_indices = [entry[0] for entry in plan]
    prs           = build_pptx_from_plan(template_path, slide_indices)
    week_days     = _days_by_week(menu_data)
    slides        = list(prs.slides)

    display_name = cfg.get('display_name', cfg['name'])
    email        = cfg.get('email', '')
    if not date_map:
        date_map = {}

    # 공휴일/선거일 사유 수집 (양식 떠있는 라벨 제거 매칭용)
    holiday_reasons = set()
    for days in week_days.values():
        for day in days:
            if day and day.get('holiday_reason'):
                holiday_reasons.add(day['holiday_reason'])

    for sidx, plan_entry in enumerate(plan):
        if sidx >= len(slides):
            break
        slide = slides[sidx]

        replace_slide_metadata(slide, display_name, email)
        replace_date_group(slide)  # 날짜 그룹 텍스트 비우기 (날짜는 테이블 셀 안에 삽입)
        remove_holiday_labels(slide, holiday_reasons)  # 양식 떠있는 공휴일 라벨 제거
        replace_origin_text(slide, origin_text)
        replace_material_text(slide, material_text)

        _, sections, opts = plan_entry
        if opts.get('skip') or not sections:
            continue

        table = _get_table(slide)
        if table is None:
            continue

        _render_slide(table, plan_entry, week_days, cfg)

        # 공휴일 셀 위에 떠서 날짜를 가리는 장식 그림(선거도장 등) 제거
        remove_holiday_cover_pics(slide, sections, week_days)

    _fix_allergy_only(prs)
    _check_box_overlap(prs, branch_label=cfg.get('name', out_path))
    prs.save(out_path)
    return out_path
