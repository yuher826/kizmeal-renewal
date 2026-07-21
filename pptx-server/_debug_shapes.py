# -*- coding: utf-8 -*-
import sys
import os

# Windows 콘솔 UTF-8 강제
sys.stdout.reconfigure(encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(BASE_DIR, "templates", "2026년 6월 식단표(양식).pptx")

# 파일 존재 확인
if not os.path.exists(TEMPLATE_PATH):
    # glob fallback
    import glob
    candidates = glob.glob(os.path.join(BASE_DIR, "templates", "*.pptx"))
    yangsik = [p for p in candidates if "양식" in p]
    if yangsik:
        TEMPLATE_PATH = yangsik[0]
    else:
        print("ERROR: template not found")
        sys.exit(1)

print("TEMPLATE:", TEMPLATE_PATH)

prs = Presentation(TEMPLATE_PATH)
slide = prs.slides[0]
shapes = list(slide.shapes)

print()
print("=" * 110)
print(f"  슬라이드1 도형 목록  (전체 {len(shapes)}개)")
print("=" * 110)
print(f"{'#':<4}  {'name':<28}  {'type':<18}  {'left':>10}  {'top':>10}  {'width':>10}  {'height':>10}")
print("-" * 110)

highlight_idx = []

for i, sh in enumerate(shapes):
    l = sh.left   if sh.left   is not None else 0
    t = sh.top    if sh.top    is not None else 0
    w = sh.width  if sh.width  is not None else 0
    h = sh.height if sh.height is not None else 0

    stype_name = str(sh.shape_type)
    try:
        stype_name = sh.shape_type.name
    except Exception:
        pass

    # 텍스트 포함 여부 체크
    flag = ""
    try:
        if sh.has_text_frame:
            txt = sh.text_frame.text
            if "지방선거일" in txt or "TextBox 26" in sh.name:
                flag = "  <<< HIGHLIGHT"
                highlight_idx.append(i)
    except Exception:
        pass
    if "TextBox 26" in sh.name:
        flag = "  <<< HIGHLIGHT (name)"
        if i not in highlight_idx:
            highlight_idx.append(i)

    line = f"{i:<4}  {sh.name:<28}  {stype_name:<18}  {l:>10}  {t:>10}  {w:>10}  {h:>10}{flag}"
    print(line)

print("=" * 110)

# ── 강조 도형 상세 ──────────────────────────────────────────────
if highlight_idx:
    print()
    print("[ HIGHLIGHT 도형 상세 ]")
    for i in highlight_idx:
        sh = shapes[i]
        l = sh.left or 0; t = sh.top or 0; w = sh.width or 0; h = sh.height or 0
        print(f"  index={i}  name='{sh.name}'")
        print(f"    left={l}  top={t}  width={w}  height={h}")
        try:
            if sh.has_text_frame:
                print(f"    text='{sh.text_frame.text[:80]}'")
        except Exception:
            pass
else:
    print()
    print("(지방선거일 / TextBox 26 해당 도형 없음)")

# ── col3 row0 셀 좌표 ─────────────────────────────────────────
print()
print("[ col3 / row0 셀 좌표 (날짜 3일 헤더 셀) ]")

TABLE_TYPE = 19  # MSO_SHAPE_TYPE.TABLE

for sh in shapes:
    if sh.shape_type == TABLE_TYPE or (hasattr(sh, 'table')):
        try:
            tbl = sh.table
        except Exception:
            continue

        tbl_left  = sh.left  or 0
        tbl_top   = sh.top   or 0

        # col widths / row heights
        col_widths = [c.width for c in tbl.columns]
        row_heights = [r.height for r in tbl.rows]

        target_col = 3
        target_row = 0

        if target_col >= len(col_widths) or target_row >= len(row_heights):
            print(f"  table '{sh.name}': col/row 범위 초과 (cols={len(col_widths)}, rows={len(row_heights)})")
            continue

        cell_left  = tbl_left + sum(col_widths[:target_col])
        cell_top   = tbl_top  + sum(row_heights[:target_row])
        cell_width  = col_widths[target_col]
        cell_height = row_heights[target_row]

        cell_obj = tbl.cell(target_row, target_col)
        cell_text = ""
        try:
            cell_text = cell_obj.text_frame.text[:40]
        except Exception:
            pass

        print(f"  table='{sh.name}'  (table_left={tbl_left}, table_top={tbl_top})")
        print(f"  col3/row0  left={cell_left}  top={cell_top}  width={cell_width}  height={cell_height}")
        print(f"  cell text='{cell_text}'")
        print(f"  col_widths[0..5] = {col_widths[:6]}")
        print(f"  row_heights[0..2] = {row_heights[:3]}")

# ── TextBox 26 문단별 텍스트 / 폰트 상세 ────────────────────────
print()
print("[ TextBox 26 : 문단별 텍스트 + 폰트 ]")
print("(양식 템플릿 기준)")

def rgb_str(color):
    try:
        return "#{:06X}".format(color.rgb)
    except Exception:
        return str(color)

def print_tf_detail(tf, indent="  "):
    for pi, para in enumerate(tf.paragraphs):
        para_text = para.text
        print(f"{indent}[para {pi}] '{para_text}'")
        for ri, run in enumerate(para.runs):
            font = run.font
            size  = font.size.pt if font.size else None
            bold  = font.bold
            color = rgb_str(font.color.rgb) if (font.color and font.color.type is not None) else "inherited"
            print(f"{indent}  run {ri}: '{run.text}'  size={size}pt  bold={bold}  color={color}")

tb26 = next((s for s in shapes if s.name == "TextBox 26"), None)
if tb26 and tb26.has_text_frame:
    print_tf_detail(tb26.text_frame)
else:
    print("  TextBox 26 없음")

# ── 강동E 생성 결과물 col3/row0 셀 텍스트 상세 ─────────────────
print()
print("[ 강동E 생성 결과물 : MENU_TABLE col3/row0 셀 텍스트 ]")

OUTPUT_PATH = os.path.join(BASE_DIR, "output", "강동E_202606.pptx")
if not os.path.exists(OUTPUT_PATH):
    # fallback: 가장 최근 강동E 출력물
    import glob
    cands = sorted(glob.glob(os.path.join(BASE_DIR, "output", "강동E*.pptx")))
    OUTPUT_PATH = cands[-1] if cands else None

if OUTPUT_PATH and os.path.exists(OUTPUT_PATH):
    print(f"  파일: {OUTPUT_PATH}")
    prs2 = Presentation(OUTPUT_PATH)
    slide2 = prs2.slides[0]
    for sh in slide2.shapes:
        if sh.name == "MENU_TABLE":
            try:
                tbl2 = sh.table
            except Exception:
                continue
            col_widths2  = [c.width for c in tbl2.columns]
            row_heights2 = [r.height for r in tbl2.rows]
            cell2 = tbl2.cell(0, 3)
            print(f"  col3/row0 전체 텍스트: '{cell2.text_frame.text}'")
            print(f"  문단별 상세:")
            print_tf_detail(cell2.text_frame, indent="    ")
            break
else:
    print(f"  출력 파일 없음: {OUTPUT_PATH}")
