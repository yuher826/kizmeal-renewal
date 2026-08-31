# -*- coding: utf-8 -*-
"""
_init_table_cells() 중식/간식 template 추출 회귀 테스트 — pytest 불필요,
직접 실행. 실행: python test_init_table_cells.py

실제 python-pptx 표를 즉석에서 만들어 pptx_generator._init_table_cells()를
그대로 호출한다(재구현이 아니라 실제 함수를 검증) — col1(월요일) 고정이
아니라 "데이터 열 어디에 있든" 동작해야 한다는 게 이번 수정의 핵심이라,
가짜/모의 로직으로는 이 버그를 재현·검증할 수 없다.
"""
import sys
from pptx import Presentation
from pptx.util import Inches

import pptx_generator as pg

N_DATA_COLS = 5  # col1~col5 (월~금). 9월 실물 양식과 동일 열 수(라벨 1 + 5)


def make_table(n_rows, n_cols):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    gframe = slide.shapes.add_table(n_rows, n_cols, Inches(1), Inches(1), Inches(6), Inches(3))
    return gframe.table


def run_count(table, ri, ci):
    tc = table.cell(ri, ci)._tc
    txBody = tc.find(f'{{{pg._NS_A}}}txBody')
    if txBody is None:
        return 0
    return len(txBody.findall(f'.//{{{pg._NS_A}}}r'))


def case_lunch_extraction(name, data_col_with_run):
    """
    lunch 행 1개짜리 표에서, data_col_with_run 번째 데이터 열(1-base)에만
    placeholder 텍스트를 넣고 나머지는 빈 채로 _init_table_cells를 호출한다.
    ★col1이 비어 있어도 다른 열에 데이터가 있으면 lunch_tmpl을 찾아 나머지
    빈 열까지 전부 채워야 한다(수정 전엔 col1이 비면 실패했음).
    """
    table = make_table(1, N_DATA_COLS + 1)
    table.cell(0, data_col_with_run).text = '흰쌀밥'

    pg._init_table_cells(table, ['lunch'])

    filled = [run_count(table, 0, ci) > 0 for ci in range(1, N_DATA_COLS + 1)]
    ok = all(filled)
    status = 'PASS' if ok else 'FAIL'
    print(f'[{status}] {name} — 데이터가 있던 열=col{data_col_with_run}, '
          f'적용 후 채워진 열={[f"col{i+1}" for i, f in enumerate(filled) if f]}')
    return ok


def case_snack_independent_of_lunch():
    """lunch 행엔 데이터가 아예 없고 snack 행에만 있어도 snack_tmpl은
    독립적으로 찾아야 한다(어느 한쪽이 없다고 다른 쪽까지 실패하면 안 됨)."""
    table = make_table(2, N_DATA_COLS + 1)  # row0=lunch(전부 빈칸), row1=morning_snack
    table.cell(1, 2).text = '사과'  # col2에만 간식 데이터

    pg._init_table_cells(table, ['lunch', 'morning_snack'])

    snack_filled = [run_count(table, 1, ci) > 0 for ci in range(1, N_DATA_COLS + 1)]
    ok = all(snack_filled)
    status = 'PASS' if ok else 'FAIL'
    print(f'[{status}] snack_tmpl은 lunch 행이 텅 비어 있어도 독립적으로 추출됨 — '
          f'적용 후 채워진 열={[f"col{i+1}" for i, f in enumerate(snack_filled) if f]}')
    return ok


def main():
    results = []

    # ★지시된 필수 케이스★ — col1이 비고 col2에만 run이 있는 표
    results.append(case_lunch_extraction('col1 비고 col2에만 데이터 (9월 케이스, 1일=화)', 2))

    # ★9월 하나로 끝내지 않기★ — 데이터가 col1~col5 중 어디에 있어도
    # 동작해야 한다. 열 위치는 "1일이 속한 주의 앞쪽 며칠이 전달 몫이라
    # 비는" 것과 같은 메커니즘이므로 대표 월을 주석으로 병기한다.
    results.append(case_lunch_extraction('col1에 데이터 (1일=월, 예: 2026-06)', 1))
    results.append(case_lunch_extraction('col2에 데이터 (1일=화, 예: 2026-09)', 2))
    results.append(case_lunch_extraction('col3에 데이터 (1일=수, 예: 2026-04)', 3))
    results.append(case_lunch_extraction('col4에 데이터 (1일=목, 예: 2026-10)', 4))
    results.append(case_lunch_extraction('col5에 데이터 (1일=금, 예: 2026-05)', 5))
    # 1일=토/일(예: 2026-08, 2026-11)은 1주차 전체가 전달 몫이라 디자이너가
    # 보통 다음 주(2주차)부터 정상적으로 col1을 채운다 — col1 케이스로
    # 이미 커버됨(별도 case 불필요).

    results.append(case_snack_independent_of_lunch())

    print()
    if all(results):
        print(f'✅ 전체 통과 ({len(results)}건)')
        return 0
    else:
        print(f'❌ 실패 있음 ({sum(1 for r in results if not r)}/{len(results)}건)')
        return 1


if __name__ == '__main__':
    sys.exit(main())
