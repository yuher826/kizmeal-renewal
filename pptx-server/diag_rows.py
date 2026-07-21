import os, sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from lxml import etree

NS_A   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_PML = 'http://schemas.openxmlformats.org/presentationml/2006/main'
SLIDE_H = 10_680_700

def diag(label, path):
    print(f'\n{"="*64}')
    print(f'  {label}  ({path})')
    print(f'{"="*64}')
    prs = Presentation(path)
    for si, slide in enumerate(prs.slides):
        slide_el = slide._element

        # graphicFrame(테이블) 찾기
        gf = None
        for shape in slide.shapes:
            if shape.has_table:
                gf = shape
                break
        if gf is None:
            print(f'  [slide{si+1}] 테이블 없음')
            continue

        tbl = gf.table
        table_top = gf.top

        print(f'\n  [slide{si+1}] 테이블 top={table_top:,}  shape.height={gf.height:,}')

        # 행별 높이 + 누적 top + 셀 첫 텍스트(주차 식별용)
        cum = 0
        total_row_h = 0
        for ri, row in enumerate(tbl.rows):
            rh = row.height or 0
            row_top = table_top + cum
            # 첫 번째 셀 텍스트(주차/행 레이블)
            try:
                cell_text = tbl.cell(ri, 0).text.strip().replace('\n', ' ')[:30]
            except Exception:
                cell_text = ''
            marker = ' ◀ 본식행?' if cell_text else ''
            print(f'    행[{ri:02d}] h={rh:>9,}  top={row_top:>11,}  "{cell_text}"{marker}')
            cum += rh
            total_row_h += rh

        table_bottom = table_top + total_row_h
        overflow = table_bottom - SLIDE_H

        print(f'\n  table_bottom (행합산) = {table_bottom:,}')
        print(f'  shape.height (XML ext) = {table_top + gf.height:,}')
        print(f'  슬라이드 높이          = {SLIDE_H:,}')
        if overflow > 0:
            print(f'  ⚠️  테이블이 슬라이드를 {overflow:,} EMU 초과')
        else:
            print(f'  ✅ 테이블이 슬라이드 안에 있음 (여유 {-overflow:,} EMU)')

BASE = os.path.dirname(os.path.abspath(__file__))
diag('15행 (type=C / S1)', os.path.join(BASE, 'test_15row.pptx'))
diag('10행 (type=A / S2)', os.path.join(BASE, 'test_10row.pptx'))
