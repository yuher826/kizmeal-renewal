import os, sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
TEMPLATE = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        'templates', '2026년 6월 식단표(양식).pptx')

TARGET = {'ORIGIN_BOX', 'MATERIAL_BOX'}

def iter_all(shapes):
    for s in shapes:
        yield s
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_all(s.shapes)

prs = Presentation(TEMPLATE)
for si, slide in enumerate(prs.slides):
    print(f'\n[slide{si+1}]')
    for shape in iter_all(slide.shapes):
        if shape.name not in TARGET:
            continue
        sp_el = shape._element
        txBody = sp_el.find(f'.//{{{NS_A}}}txBody')
        bodyPr = txBody.find(f'{{{NS_A}}}bodyPr') if txBody is not None else None

        cy = shape.height
        top = shape.top
        print(f'  {shape.name}: top={top:,}  cy={cy:,}  bottom={top+cy:,}')

        if bodyPr is not None:
            attrs = dict(bodyPr.attrib)
            children = [etree.QName(c).localname for c in bodyPr]
            print(f'    bodyPr attrs   : {attrs}')
            print(f'    bodyPr children: {children}')
        else:
            print(f'    bodyPr: 없음')

# 10행(S2) 주요 좌표 계산
print('\n--- 10행(S2) 공간 계산 ---')
slide2 = prs.slides[1]
coords = {}
for shape in iter_all(slide2.shapes):
    if shape.name in TARGET:
        coords[shape.name] = {'top': shape.top, 'cy': shape.height}

if 'ORIGIN_BOX' in coords and 'MATERIAL_BOX' in coords:
    o = coords['ORIGIN_BOX']
    m = coords['MATERIAL_BOX']
    available = m['top'] - 8_550_000   # ORIGIN_Y=8550000 기준 사용 가능 공간
    print(f'  ORIGIN_Y(설정값)     = 8,550,000')
    print(f'  ORIGIN_BOX 원본 cy   = {o["cy"]:,}')
    print(f'  MATERIAL_BOX top     = {m["top"]:,}')
    print(f'  사용 가능 공간       = {available:,}  (MATERIAL top - ORIGIN_Y)')
    print(f'  초과분               = {o["cy"] - available:,}  (cy - 사용가능)')
    print(f'  → 5.5pt 1줄 ≈ 200,000 EMU 기준: cy로 수용 가능 줄 수 ≈ {o["cy"]//200000:.1f}줄')
