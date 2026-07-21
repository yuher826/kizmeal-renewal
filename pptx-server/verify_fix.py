import os, sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx_generator import generate

TEMPLATE = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        'templates', '2026년 6월 식단표(양식).pptx')
OUT_15 = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'test_15row_v2.pptx')
OUT_10 = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'test_10row_v2.pptx')

EMPTY_MENU = {'weeks': {str(w): {'days': [None] * 5} for w in range(1, 6)}}

def cfg(name, type_code):
    return {'name': name, 'display_name': name, 'type': type_code, 'email': ''}

print('=== 15행 생성 (type=C, S1) ===')
ORIGIN = {
    'body': (
        '쌀-현미·흑미·찹쌀(국내산 친환경), 소고기(국내산 육우), 돼지고기(국내산), '
        '닭고기(국내산), 달걀(국내산 친환경), 두부(국내산 친환경 콩), '
        '시금치·당근·양파·감자·고추(국내산 친환경), 김치(국내산 친환경), '
        '멸치(국내산), 미역(국내산), 된장·간장·고추장(국내산 친환경), '
        '참기름(국내산), 식용유(수입산 대두), 밀가루(수입산)'
    ),
    'disclaimer': '※ 원산지는 변동될 수 있으며, 변경 시 별도 공지합니다.',
}
MATERIAL = (
    '달걀①, 우유②, 밀③, 땅콩⑤, 대두⑥, 밀③, 고등어⑧, 게⑨, '
    '새우⑩, 돼지고기⑪, 복숭아⑫, 토마토⑬, 아황산류⑭, 호두⑮, '
    '닭고기⑯, 쇠고기⑰, 오징어⑱, 조개류⑲'
)

generate(cfg('테스트15행', 'C'), EMPTY_MENU, TEMPLATE, OUT_15,
         origin_text=ORIGIN, material_text=MATERIAL)

print('\n=== 10행 생성 (type=A, S2) ===')
generate(cfg('테스트10행', 'A'), EMPTY_MENU, TEMPLATE, OUT_10,
         origin_text=ORIGIN, material_text=MATERIAL)

# ── 좌표 출력 ──────────────────────────────────────────────────────
TARGET = {
    'ALLERGY_BOX', 'ORIGIN_BOX', 'MATERIAL_BOX',
    '그래픽 118',                        # 로고
    '그림 3', '그림 224', '그림 193',    # 클립보드 (슬라이드별)
}
SLIDE_H = 10_680_700

# 양식 원본 좌표 (비교 기준)
ORIGINAL = {
    'ORIGIN_BOX':   {'top': 8_395_899},
    'MATERIAL_BOX': {'top': 9_316_195},
    '그림 3':       {'top': 8_422_882},   # slide1 클립보드
    '그림 224':     {'top': 8_422_882},   # slide2 클립보드
    '그래픽 118':   {'top': 9_900_789},
}

def iter_all(shapes):
    for s in shapes:
        yield s
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_all(s.shapes)

def print_coords(label, path):
    print(f'\n{"="*60}')
    print(f'  {label}')
    print(f'{"="*60}')
    prs = Presentation(path)
    for si, slide in enumerate(prs.slides):
        print(f'  [slide{si+1}]')
        for shape in iter_all(slide.shapes):
            if shape.name not in TARGET:
                continue
            orig = ORIGINAL.get(shape.name, {}).get('top')
            moved = '' if orig is None else (
                '  ← 원본유지 ✅' if shape.top == orig else
                f'  ← 이동됨 (원본:{orig:,})'
            )
            bottom = shape.top + shape.height
            overflow = '  ⚠️ 슬라이드 밖!' if bottom > SLIDE_H else ''
            print(f'    {shape.name:20s}  top={shape.top:>11,}  bottom={bottom:>11,}  cy={shape.height:>10,}{moved}{overflow}')

print_coords('15행 결과 (ORIGIN/MATERIAL/클립보드/로고 → 원본 좌표 기대)', OUT_15)
print_coords('10행 결과 (ALLERGY/ORIGIN/클립보드/로고 → 어제 완성본 기대)', OUT_10)
