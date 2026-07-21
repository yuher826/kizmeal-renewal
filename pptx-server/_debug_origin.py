# -*- coding: utf-8 -*-
"""원본_강동ECC_영양사.pptx 분석: 3일(선거일) 날짜 표시 방식 확인"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8', errors='replace')

from pptx import Presentation
from lxml import etree

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
ORIGIN_PATH = os.path.join(BASE_DIR, "원본_강동ECC_영양사.pptx")

if not os.path.exists(ORIGIN_PATH):
    import glob
    cands = glob.glob(os.path.join(BASE_DIR, "*영양사*.pptx"))
    if cands:
        ORIGIN_PATH = cands[0]
    else:
        print("ERROR: 원본 파일 없음")
        sys.exit(1)

print("ORIGIN:", ORIGIN_PATH)

prs = Presentation(ORIGIN_PATH)
slide = prs.slides[0]
shapes = list(slide.shapes)

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def run_color_hex(run):
    """run 폰트색 hex 추출. 실패 시 'inherited'."""
    try:
        c = run.font.color
        if c is not None and c.type is not None:
            return "#{:06X}".format(c.rgb)
    except Exception:
        pass
    # rPr solidFill 직접 파싱 fallback
    try:
        rPr = run._r.find(f'{{{_NS_A}}}rPr')
        if rPr is not None:
            srgb = rPr.find(f'.//{{{_NS_A}}}srgbClr')
            if srgb is not None:
                return "#" + srgb.get('val').upper()
    except Exception:
        pass
    return "inherited"


def run_size_pt(run):
    try:
        return run.font.size.pt if run.font.size else None
    except Exception:
        return None


def dump_tf(tf, indent="    "):
    for pi, para in enumerate(tf.paragraphs):
        print(f"{indent}[para {pi}] '{para.text}'")
        for ri, run in enumerate(para.runs):
            print(f"{indent}  run {ri}: '{run.text}'  "
                  f"size={run_size_pt(run)}pt  color={run_color_hex(run)}")


# ════════════════════════════════════════════════════════════
# 1. MENU_TABLE col3/row0 셀
# ════════════════════════════════════════════════════════════
print()
print("=" * 90)
print("[ 1. 원본 MENU_TABLE col3/row0 (3일 헤더 셀) ]")
print("=" * 90)

menu_table = None
for sh in shapes:
    try:
        if sh.name == "MENU_TABLE" and sh.has_table:
            menu_table = sh
            break
    except Exception:
        continue

if menu_table is None:
    # 폴백: 첫 table
    for sh in shapes:
        try:
            if sh.has_table:
                menu_table = sh
                print(f"  (MENU_TABLE 이름 없음 → '{sh.name}' 사용)")
                break
        except Exception:
            continue

if menu_table is not None:
    tbl = menu_table.table
    cell = tbl.cell(0, 3)
    full = cell.text_frame.text
    print(f"  table='{menu_table.name}'")
    print(f"  col3/row0 전체 텍스트: {repr(full)}")
    if full.strip():
        print("  → 셀에 텍스트 있음. 문단별:")
        dump_tf(cell.text_frame)
    else:
        print("  → 셀 비어 있음 (텍스트 없음)")
else:
    print("  표를 찾지 못함")


# ════════════════════════════════════════════════════════════
# 2. "3" 또는 "지방선거일/선거일" 텍스트 가진 모든 도형
# ════════════════════════════════════════════════════════════
print()
print("=" * 90)
print("[ 2. '3' 또는 '선거일' 텍스트 포함 도형 전수 조사 (그룹 내부 포함) ]")
print("=" * 90)


def iter_all(shapes):
    for s in shapes:
        yield s
        try:
            if s.shape_type == 6:  # GROUP
                yield from iter_all(s.shapes)
        except Exception:
            pass


def shape_coords(sh):
    l = sh.left if sh.left is not None else None
    t = sh.top if sh.top is not None else None
    w = sh.width if sh.width is not None else None
    h = sh.height if sh.height is not None else None
    return l, t, w, h


targets = ["3", "선거일", "지방선거일"]
found_any = False

for sh in iter_all(shapes):
    try:
        if not sh.has_text_frame:
            continue
    except Exception:
        continue
    txt = sh.text_frame.text
    txt_stripped = txt.strip()
    # 정확히 "3" 이거나 선거일 포함
    hit = (txt_stripped == "3") or ("선거일" in txt)
    if not hit:
        continue
    found_any = True
    l, t, w, h = shape_coords(sh)
    stype = sh.shape_type
    try:
        stype = sh.shape_type.name
    except Exception:
        pass
    print(f"  ▸ name='{sh.name}'  type={stype}")
    print(f"    text={repr(txt)}")
    print(f"    left={l} top={t} width={w} height={h}")
    dump_tf(sh.text_frame, indent="    ")
    print()

if not found_any:
    print("  (표 셀 밖에서 '3'/'선거일' 독립 도형 없음)")


# ════════════════════════════════════════════════════════════
# 2-b. 표 셀 전체에서 "3"/"선거일" 검색 (셀은 도형 순회에 안 잡힘)
# ════════════════════════════════════════════════════════════
print("-" * 90)
print("[ 2-b. MENU_TABLE 모든 셀 중 '3'/'선거일' 포함 셀 ]")
if menu_table is not None:
    tbl = menu_table.table
    n_rows = len(tbl.rows)
    n_cols = len(tbl.columns)
    for r in range(n_rows):
        for c in range(n_cols):
            cell = tbl.cell(r, c)
            ctext = cell.text_frame.text
            cs = ctext.strip()
            if cs == "3" or "선거일" in ctext:
                print(f"  cell(row={r}, col={c}) text={repr(ctext)}")
                dump_tf(cell.text_frame, indent="      ")


# ════════════════════════════════════════════════════════════
# 3. 결론용 요약
# ════════════════════════════════════════════════════════════
print()
print("=" * 90)
print("[ 3. 결론 요약 ]")
print("=" * 90)
if menu_table is not None:
    cell = menu_table.table.cell(0, 3)
    full = cell.text_frame.text
    if full.strip():
        runs = [r for p in cell.text_frame.paragraphs for r in p.runs]
        colors = [run_color_hex(r) for r in runs]
        print(f"  ▶ col3/row0 셀에 텍스트 존재: {repr(full)}")
        print(f"  ▶ 셀 내 run 색상들: {colors}")
    else:
        print("  ▶ col3/row0 셀 비어 있음 → 날짜는 셀 밖 독립 도형 가능성")
