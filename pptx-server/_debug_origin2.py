# -*- coding: utf-8 -*-
"""원본 표 전체 셀 덤프 — 날짜가 어디에 어떻게 들어있는지 확인"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8', errors='replace')
from pptx import Presentation

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
BASE = os.path.dirname(os.path.abspath(__file__))
path = os.path.join(BASE, "원본_강동ECC_영양사.pptx")

prs = Presentation(path)
slide = prs.slides[0]


def color_hex(run):
    try:
        c = run.font.color
        if c is not None and c.type is not None:
            return "#{:06X}".format(c.rgb)
    except Exception:
        pass
    try:
        rPr = run._r.find(f'{{{_NS_A}}}rPr')
        srgb = rPr.find(f'.//{{{_NS_A}}}srgbClr') if rPr is not None else None
        if srgb is not None:
            return "#" + srgb.get('val').upper()
    except Exception:
        pass
    return "inh"


tbl = None
for sh in slide.shapes:
    try:
        if sh.has_table:
            tbl = sh.table
            print(f"TABLE: name='{sh.name}'  rows={len(tbl.rows)} cols={len(tbl.columns)}")
            break
    except Exception:
        pass

# row0 (헤더/날짜 행) 전체 + row1~2 일부
for r in range(min(3, len(tbl.rows))):
    print(f"\n--- row {r} ---")
    for c in range(len(tbl.columns)):
        cell = tbl.cell(r, c)
        txt = cell.text_frame.text
        runs = [run for p in cell.text_frame.paragraphs for run in p.runs]
        cols = [(run.text, color_hex(run),
                 (run.font.size.pt if run.font.size else None)) for run in runs]
        print(f"  col{c}: text={repr(txt)}")
        for rt, ch, sz in cols:
            print(f"        run '{rt}'  {ch}  {sz}pt")
