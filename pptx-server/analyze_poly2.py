from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def iter_all(shapes):
    for s in shapes:
        yield s
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_all(s.shapes)

prs = Presentation("templates/2026년 6월 식단표(양식).pptx")
slide = prs.slides[0]

for shape in iter_all(slide.shapes):
    if shape.name in ["ALLERGY_BOX", "그림 224", "그래픽 118"]:
        print(f"{shape.name}: top={shape.top}, left={shape.left}, w={shape.width}, h={shape.height}")
