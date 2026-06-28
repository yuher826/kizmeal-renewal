# -*- coding: utf-8 -*-
"""bodyPr 비교 분석 — 3일(공휴일) vs 10일(평일) vs 1일(정상 메뉴)"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from lxml import etree
from pptx import Presentation

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
path = 'output/강동E_202606.pptx'

prs = Presentation(path)
slide = list(prs.slides)[0]
table = None
for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        break


def dump_bodyPr(label, ri, ci):
    cell   = table.cell(ri, ci)
    txBody = cell._tc.find(f'{{{NS_A}}}txBody')
    bodyPr = txBody.find(f'{{{NS_A}}}bodyPr') if txBody is not None else None

    print(f'── {label} [row{ri}, col{ci}] ──')
    if bodyPr is None:
        print('  bodyPr: 없음')
    else:
        attrs    = dict(bodyPr.attrib)
        children = list(bodyPr)
        print(f'  bodyPr 속성: {attrs if attrs else "(없음 — 완전히 빈 태그)"}')
        print(f'  bodyPr 자식: {[c.tag.split("}")[1] for c in children] if children else "(없음)"}')
        print(f'  bodyPr XML: {etree.tostring(bodyPr, pretty_print=False).decode()}')

    tcPr   = cell._tc.find(f'{{{NS_A}}}tcPr')
    anchor = tcPr.get('anchor') if tcPr is not None else None
    print(f'  tcPr anchor: {anchor}')

    if txBody is not None:
        paras = txBody.findall(f'{{{NS_A}}}p')
        print(f'  단락 수: {len(paras)}')
        for pi, p in enumerate(paras):
            pPr  = p.find(f'{{{NS_A}}}pPr')
            runs = p.findall(f'{{{NS_A}}}r')
            algn = pPr.get('algn') if pPr is not None else None
            texts = [r.find(f'{{{NS_A}}}t').text for r in runs
                     if r.find(f'{{{NS_A}}}t') is not None]
            print(f'    단락[{pi}] algn={algn}, run수={len(runs)}, text={texts}')
    print()


# 1일 본식(정상 메뉴)
dump_bodyPr('1일 본식 [정상 메뉴]', 0, 1)

# 3일 본식(공휴일 — 날짜 안 보임)
dump_bodyPr('3일 본식 [공휴일]', 0, 3)

# 10일 본식(평일)
dump_bodyPr('10일 본식 [평일]', 3, 3)

# ── clear_cell / _set_date_only_cell 소스 확인 ───────────────────────
print('=' * 60)
print('clear_cell 소스')
print('=' * 60)
import inspect, pptx_generator
print(inspect.getsource(pptx_generator.clear_cell))

print('=' * 60)
print('_set_date_only_cell 소스')
print('=' * 60)
print(inspect.getsource(pptx_generator._set_date_only_cell))
