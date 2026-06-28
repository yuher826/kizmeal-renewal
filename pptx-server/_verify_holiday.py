# -*- coding: utf-8 -*-
"""공휴일 날짜+사유 표시 검증: 강동E 3일 + 평일 회귀 + holiday_reason 파싱"""
import sys, re
sys.stdout.reconfigure(encoding='utf-8')
from lxml import etree
from pptx import Presentation
import glob

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# ── holiday_reason 파싱 검증 ─────────────────────────────────────────
print('=== parse_header holiday_reason 확인 ===')
from read_excel import parse_header
cases = [
    '수  6/3\n공휴일\n(선거일)',
    '월  6/1',
    '수  6/3\n공휴일',           # 괄호 없는 공휴일 (엣지케이스)
    '수  6/3\n공휴일\n(지방선거일)',
]
for s in cases:
    r = parse_header(s)
    print(f'  입력: {repr(s[:30])}')
    print(f'  → is_holiday={r["is_holiday"]}, holiday_reason={repr(r["holiday_reason"])}')
    print()

# ── 강동E PPTX 셀 덤프 ───────────────────────────────────────────────
candidates = glob.glob('output/*강동E*.pptx') + glob.glob('output/*강동e*.pptx')
if not candidates:
    print('강동E 파일을 output/ 에서 찾지 못했습니다.')
    sys.exit(1)
path = candidates[0]
print(f'=== 강동E PPTX: {path} ===\n')

prs = Presentation(path)
slide = list(prs.slides)[0]
table = None
for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        break
if table is None:
    print('테이블 없음'); sys.exit(1)

print(f'테이블: {len(table.rows)}행 × {len(table.columns)}열\n')


def dump_cell(label, ri, ci):
    cell = table.cell(ri, ci)
    tc = cell._tc
    txBody = tc.find(f'{{{NS_A}}}txBody')

    paragraphs = []
    if txBody is not None:
        for p in txBody.findall(f'{{{NS_A}}}p'):
            runs = p.findall(f'{{{NS_A}}}r')
            text = ''.join(
                (r.find(f'{{{NS_A}}}t').text or '')
                for r in runs if r.find(f'{{{NS_A}}}t') is not None
            )
            # 이 단락 첫 run의 색상 추출
            color = None
            if runs:
                rPr = runs[0].find(f'{{{NS_A}}}rPr')
                if rPr is not None:
                    sf = rPr.find(f'.//{{{NS_A}}}solidFill')
                    if sf is not None:
                        sc = sf.find(f'{{{NS_A}}}srgbClr')
                        if sc is not None:
                            color = f'#{sc.get("val","?")}'
                    sz = rPr.get('sz', '?')
                    bold = rPr.get('b', '0')
            if text or color:
                paragraphs.append(f'<{text}> 색={color} sz={sz} b={bold}')
            else:
                paragraphs.append(f'<빈단락>')

    # 배경색
    tcPr = tc.find(f'{{{NS_A}}}tcPr')
    bg = '(없음)'
    if tcPr is not None:
        sf = tcPr.find(f'.//{{{NS_A}}}solidFill')
        if sf is not None:
            sc = sf.find(f'{{{NS_A}}}srgbClr')
            if sc is not None:
                bg = f'#{sc.get("val","?")}'

    print(f'  [{label}] row{ri} col{ci}  배경={bg}')
    for pg in paragraphs:
        print(f'    {pg}')
    print()


print('=== 1주차 3일(공휴일/선거일) ===')
dump_cell('본식   ', 0, 3)
dump_cell('오전간식', 1, 3)
dump_cell('오후간식', 2, 3)

print('=== 1주차 1일(평일, 회귀 확인) ===')
dump_cell('본식   ', 0, 1)
dump_cell('오전간식', 1, 1)

print('=== 2주차 10일(평일, 포맷 확인) ===')
dump_cell('본식   ', 3, 3)
dump_cell('오전간식', 4, 3)
