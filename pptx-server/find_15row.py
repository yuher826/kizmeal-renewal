from pptx import Presentation
import os

folder = os.path.join(os.path.dirname(os.path.abspath(__file__)), "templates")
for fname in sorted(os.listdir(folder)):
    if fname.endswith(".pptx"):
        try:
            prs = Presentation(os.path.join(folder, fname))
            slide = prs.slides[0]
            for shape in slide.shapes:
                if shape.has_table:
                    rows = len(shape.table.rows)
                    mark = "← 15행" if rows >= 15 else ""
                    print(f"{rows}행  {fname}  {mark}")
                    break
        except Exception as e:
            print(f"ERROR {fname}: {e}")
