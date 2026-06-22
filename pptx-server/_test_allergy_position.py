# -*- coding: utf-8 -*-
"""
송도POLY 1개만 생성해서 알레르기 박스 위치를 확인하는 로컬 테스트.
실행:
  $env:YEAR=2026; $env:MONTH=6; python pptx-server/_test_allergy_position.py
"""
import os, sys, tempfile
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

# .env.local 자동 로드 (python-dotenv 있을 때만)
try:
    from dotenv import load_dotenv
    load_dotenv(os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', '.env.local'))
except ImportError:
    pass

SUPABASE_URL = os.environ.get('SUPABASE_URL', '')
SUPABASE_KEY = os.environ.get('SUPABASE_SERVICE_KEY', '')
YEAR         = int(os.environ.get('YEAR', 2026))
MONTH        = int(os.environ.get('MONTH', 6))

if not SUPABASE_URL or not SUPABASE_KEY:
    print('ERROR: SUPABASE_URL / SUPABASE_SERVICE_KEY 환경변수 없음')
    sys.exit(1)

from supabase_uploader import SupabaseREST
from pptx_generator import generate as gen_pptx, _fix_allergy_position
from pptx import Presentation
from pptx.util import Emu
from app_actions import (
    _adapt_ts_menu_data, _build_date_map,
    fetch_menu_data, fetch_branch_cfgs, fetch_diet_settings,
)

TARGET = '송도POLY'

client = SupabaseREST(SUPABASE_URL, SUPABASE_KEY)

# 전역 함수에 client 바인딩
import app_actions
app_actions.client = client
app_actions.YEAR   = YEAR
app_actions.MONTH  = MONTH

print(f'[테스트] {YEAR}년 {MONTH}월 / 대상: {TARGET}')

raw_menu = fetch_menu_data()
adapted  = _adapt_ts_menu_data(raw_menu)
date_map = _build_date_map(adapted)

settings = fetch_diet_settings()
if settings.get('origin_body'):
    origin_text = {'body': settings['origin_body'], 'disclaimer': settings.get('origin_disclaimer', '')}
else:
    origin_text = None
material_text = settings.get('material_text')

branch_cfgs = fetch_branch_cfgs()
cfg = next((c for c in branch_cfgs if c['name'] == TARGET), None)
if cfg is None:
    print(f'ERROR: {TARGET} 브랜치를 찾지 못했습니다. 전체 목록:')
    for c in branch_cfgs:
        print(' -', c['name'])
    sys.exit(1)

TEMPLATE_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'templates', '2026년 6월 식단표(양식).pptx')

with tempfile.TemporaryDirectory() as tmp:
    out_path = os.path.join(tmp, 'test_allergy.pptx')
    gen_pptx(cfg, adapted, TEMPLATE_PATH, out_path, date_map=date_map,
             origin_text=origin_text, material_text=material_text)

    # 생성된 파일을 다시 열어서 위치 검증
    prs = Presentation(out_path)
    slide = prs.slides[0]

    table_shape   = None
    allergy_shape = None
    for shape in slide.shapes:
        if shape.has_table:
            table_shape = shape
        elif shape.has_text_frame:
            try:
                if '알레르기' in shape.text_frame.text:
                    allergy_shape = shape
            except Exception:
                pass

    if not table_shape or not allergy_shape:
        print('⚠️  테이블 또는 알레르기 박스를 찾지 못했습니다 — 슬라이드에 해당 요소 없음')
        sys.exit(0)

    table_bottom = table_shape.top + table_shape.height
    expected_top = table_bottom + Emu(50000)
    actual_top   = allergy_shape.top

    print(f'\n── 검증 결과 ──────────────────────────────')
    print(f'  table_bottom : {table_bottom:>12,} EMU')
    print(f'  expected_top : {expected_top:>12,} EMU  (table_bottom + 50000)')
    print(f'  actual_top   : {actual_top:>12,} EMU')
    print(f'  allergy.height: {allergy_shape.height:>11,} EMU')

    if actual_top == expected_top:
        print('\n✅ PASS — 알레르기 박스 위치 정확')
    else:
        diff = actual_top - expected_top
        print(f'\n❌ FAIL — 차이: {diff:+,} EMU')
