# -*- coding: utf-8 -*-
"""
템플릿 이름표 자동 부여 (HANDOFF 발견① 해결)
========================================================================
디자이너 PPTX 도형에 cNvPr name(MENU_TABLE/ORIGIN_BOX/MATERIAL_BOX/ALLERGY_BOX)을
자동으로 부여한다.

★ 탐색 기준은 도형 이름이 아니라 텍스트다. 도형 이름은 슬라이드마다·달마다 바뀌지만
  텍스트는 1~8월 실측 결과 내내 불변임이 확인됨(HANDOFF "식단표 자동화" 조사 발견④).

★ python-pptx(lxml) 기반으로만 작성한다. `_apply_template_names.py`(ElementTree
  전체 재작성 방식)는 p: 접두사가 날아가고 MS 확장 네임스페이스가 ns2:/ns3:/ns5:로
  개명되며 파일이 수천 바이트 단위로 줄어드는 부작용이 실측 확인됨 — 이식 금지.
  python-pptx로 동일 작업 시 차이가 ±100바이트 이내, 네임스페이스 그대로임도 실측 확인됨.

⚠️ 일본어 금지. 모든 주석/문자열 한국어.
"""

import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

REQUIRED_NAMES = ['MENU_TABLE', 'ORIGIN_BOX', 'MATERIAL_BOX', 'ALLERGY_BOX']


def _iter_all_shapes(shapes):
    """그룹(GROUP) 내부까지 재귀 순회. ALLERGY_BOX가 그룹 안에 로고와 함께 있어 필요.
    pptx_generator.py:608 _iter_all_shapes와 동일 로직(중복 방지 위해 여기서도 보유)."""
    for s in shapes:
        yield s
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                yield from _iter_all_shapes(s.shapes)
            except Exception:
                pass


def _shape_text(shape):
    if not shape.has_text_frame:
        return ''
    return shape.text_frame.text


def _find_menu_table(shapes):
    """표(has_table) 도형 중 면적이 가장 큰 것을 MENU_TABLE 후보로 선정.
    디자이너가 실수로 작은 표를 추가로 남겼을 때 잘못 골라지는 것을 방지.
    반환: (선정된 도형 또는 None, 발견된 표 개수)"""
    candidates = []
    for s in _iter_all_shapes(shapes):
        try:
            if s.has_table:
                area = (s.width or 0) * (s.height or 0)
                candidates.append((area, s))
        except Exception:
            continue
    if not candidates:
        return None, 0
    candidates.sort(key=lambda t: t[0], reverse=True)
    return candidates[0][1], len(candidates)


def apply_names_to_slide(slide):
    """슬라이드 1장에 이름표 4개를 부여.
    이미 올바른 이름표가 붙어있으면 그대로 두고 성공으로 친다(멱등).
    반환: (부여된 이름표 리스트, 누락된 이름표 리스트, 경고 리스트)"""
    warnings = []
    all_shapes = list(_iter_all_shapes(slide.shapes))
    used = {s.name for s in all_shapes if s.name in REQUIRED_NAMES}

    # 1) MENU_TABLE — 표 구조 기준
    if 'MENU_TABLE' not in used:
        menu_table, table_count = _find_menu_table(slide.shapes)
        if table_count > 1:
            warnings.append(f'표 도형이 {table_count}개 발견됨 — 가장 큰 표를 MENU_TABLE로 선정')
        if menu_table is not None:
            menu_table.name = 'MENU_TABLE'
            used.add('MENU_TABLE')

    # 2) ORIGIN_BOX / MATERIAL_BOX / ALLERGY_BOX — 텍스트 기준
    for shape in all_shapes:
        if shape.name in REQUIRED_NAMES:
            continue  # 이미 이름표 있는 도형은 건드리지 않음
        text = _shape_text(shape)
        if not text:
            continue
        if '원산지 표기' in text:
            target = 'ORIGIN_BOX'
        elif '원재료' in text and '표시안내' in text:
            target = 'MATERIAL_BOX'
        elif '알레르기 표시' in text:
            target = 'ALLERGY_BOX'
        else:
            continue
        if target in used:
            warnings.append(f'{target} 후보 도형이 중복 발견됨 — 먼저 찾은 도형만 사용')
            continue
        shape.name = target
        used.add(target)

    applied = [n for n in REQUIRED_NAMES if n in used]
    missing = [n for n in REQUIRED_NAMES if n not in used]
    return applied, missing, warnings


def apply_template_names(pptx_bytes: bytes) -> tuple[bytes, dict]:
    """
    디자이너 PPTX의 도형에 이름표를 자동 부여한다.
    반환: (수정된 bytes, {'slide1': {'applied': [...], 'missing': [...], 'warnings': [...]}, ...})
    """
    prs = Presentation(io.BytesIO(pptx_bytes))
    report = {}
    for idx, slide in enumerate(prs.slides, start=1):
        applied, missing, warnings = apply_names_to_slide(slide)
        report[f'slide{idx}'] = {'applied': applied, 'missing': missing, 'warnings': warnings}

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue(), report


if __name__ == '__main__':
    import sys
    if len(sys.argv) < 3:
        print('사용법: python template_namer.py <입력.pptx> <출력.pptx>')
        sys.exit(1)

    with open(sys.argv[1], 'rb') as f:
        in_bytes = f.read()

    out_bytes, report = apply_template_names(in_bytes)

    with open(sys.argv[2], 'wb') as f:
        f.write(out_bytes)

    print(f'입력 {len(in_bytes):,}바이트 → 출력 {len(out_bytes):,}바이트 (차이 {len(out_bytes) - len(in_bytes):+,})')
    for slide_name, info in report.items():
        line = f"  {slide_name}: 부여={info['applied']} 누락={info['missing']}"
        if info['warnings']:
            line += f" 경고={info['warnings']}"
        print(line)
