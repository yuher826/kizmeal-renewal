from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def iter_all(s):
    for x in s:
        yield x
        if x.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_all(x.shapes)

prs = Presentation("templates/2026년 6월 식단표(양식).pptx")
for si, slide in enumerate(prs.slides):
    for shape in iter_all(slide.shapes):
        print(f"slide{si+1} [{shape.name}]: top={shape.top} left={shape.left}")
