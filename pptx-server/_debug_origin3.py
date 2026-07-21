# -*- coding: utf-8 -*-
"""원본: 평일 날짜 숫자(1~31)가 어디에 어떻게 있는지 확인
   + 양식 템플릿의 평일 날짜 처리 비교"""
import sys, os, re
sys.stdout.reconfigure(encoding='utf-8', errors='replace')
from pptx import Presentation

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
BASE = os.path.dirname(os.path.abspath(__file__))

DATE_RE = re.compile(r'^(0?[1-9]|[12][0-9]|3[01])$')


def color_hex_run(run):
    try:
        c = run.font.color
        if c is not None and c.type is not None:
            return "#{:06X}".format(c.rgb)
    except Exception:
        pass
    try:
        rPr = run._r.find(f'{{{_NS_A}}}rPr')
        srgb = rPr.find(f'.//{{{_NS_A}}}srgbClr') if rPr is not None else None
        if srgb is not None:
            return "#" + srgb.get('val').upper()
    except Exception:
        pass
    return "inh"


def analyze(path, label):
    print("=" * 80)
    print(f"[ {label} ]  {os.path.basename(path)}")
    print("=" * 80)
    prs = Presentation(path)
    slide = prs.slides[0]
    spTree = slide.shapes._spTree

    # ── (A) 떠있는 sp 도형 중 날짜 숫자 ──────────────────────
    print("\n(A) 표 밖/그룹 내 '날짜 숫자' sp 도형:")
    found = []
    for sp in spTree.findall(f'.//{{{_NS_P}}}sp'):
        txt = ''.join(t.text or '' for t in sp.findall(f'.//{{{_NS_A}}}t')).strip()
        if DATE_RE.match(txt):
            off = sp.find(f'.//{{{_NS_A}}}off')
            x = off.get('x') if off is not None else '?'
            y = off.get('y') if off is not None else '?'
            srgb = sp.find(f'.//{{{_NS_A}}}srgbClr')
            col = "#" + srgb.get('val').upper() if srgb is not None else "inh"
            sz_el = sp.find(f'.//{{{_NS_A}}}rPr')
            sz = sz_el.get('sz') if sz_el is not None else '?'
            nv = sp.find(f'.//{{{_NS_P}}}cNvPr')
            nm = nv.get('name') if nv is not None else '?'
            found.append((int(txt), txt, nm, x, y, col, sz))
    found.sort()
    if found:
        print(f"    {'date':>4}  {'name':<20}  {'x':>9}  {'y':>9}  {'color':<9}  sz")
        for d, txt, nm, x, y, col, sz in found:
            print(f"    {txt:>4}  {nm:<20}  {x:>9}  {y:>9}  {col:<9}  {sz}")
        print(f"    → 날짜 숫자: {sorted(d for d,*_ in found)}")
    else:
        print("    (없음)")

    # ── (B) 표 셀 안 날짜 prefix 여부 ────────────────────────
    print("\n(B) 표 점심행(row0) 셀이 날짜 숫자로 시작하나?:")
    tbl = None
    for sh in slide.shapes:
        try:
            if sh.has_table:
                tbl = sh.table
                break
        except Exception:
            pass
    if tbl is None:
        print("    (표 없음)")
        return
    for c in range(len(tbl.columns)):
        cell = tbl.cell(0, c)
        paras = cell.text_frame.paragraphs
        first = paras[0].text if paras else ''
        # 첫 문단 색/크기
        first_runs = paras[0].runs if paras else []
        meta = ""
        if first_runs:
            r0 = first_runs[0]
            meta = f"  [color={color_hex_run(r0)} size={r0.font.size.pt if r0.font.size else None}pt]"
        is_date = DATE_RE.match(first.strip()) is not None
        mark = " <== 날짜숫자 문단" if is_date else ""
        print(f"    col{c} 첫문단: {repr(first)}{meta}{mark}")


def date_colors(path):
    """떠있는 날짜 숫자 sp 의 {날짜int: color} 반환 + 색별 집계."""
    prs = Presentation(path)
    slide = prs.slides[0]
    spTree = slide.shapes._spTree
    out = {}
    for sp in spTree.findall(f'.//{{{_NS_P}}}sp'):
        txt = ''.join(t.text or '' for t in sp.findall(f'.//{{{_NS_A}}}t')).strip()
        if DATE_RE.match(txt):
            srgb = sp.find(f'.//{{{_NS_A}}}srgbClr')
            col = "#" + srgb.get('val').upper() if srgb is not None else "inh"
            out[int(txt)] = col
    return out


# ── 1. 원본 6일(현충일) "06" 색 ──────────────────────────────
origin = os.path.join(BASE, "원본_강동ECC_영양사.pptx")
analyze(origin, "원본 영양사 정답본")

print()
print("#" * 80)
print("[ 공휴일 빨강 사례 조사 ]")
print("#" * 80)

print("\n(1) 원본_강동ECC 6일(현충일) '06' 날짜 색:")
oc = date_colors(origin)
if 6 in oc:
    print(f"    06 색 = {oc[6]}")
else:
    print("    '06' 떠있는 날짜 도형 없음 (6/6은 토요일·주말이라 날짜칸 자체 부재 가능)")
print(f"    참고 3일(선거일) 색 = {oc.get(3, '없음')}")
print(f"    전체 날짜별 색: {dict(sorted(oc.items()))}")

# ── 2. 다른 영양사 원본 PPTX 검색 ───────────────────────────
print("\n(2) pptx-server 하위 '영양사/원본' PPTX 검색:")
import glob
patterns = ["**/*영양사*.pptx", "**/원본*.pptx"]
files = set()
for pat in patterns:
    for f in glob.glob(os.path.join(BASE, pat), recursive=True):
        files.add(f)
files = sorted(files)
if files:
    for f in files:
        print(f"    - {os.path.relpath(f, BASE)}")
else:
    print("    (없음)")

# ── 3. 각 원본의 색 사례 종합 ───────────────────────────────
print("\n(3) 각 원본 날짜색 종합 (빨강 계열 사용 여부):")
RED_PREFIXES = ("#C", "#F", "#E", "#D", "#B", "#A")  # 대략 빨강/주황 계열 후보
for f in files:
    try:
        dc = date_colors(f)
    except Exception as e:
        print(f"    {os.path.basename(f)}: 분석 실패 {e}")
        continue
    colorset = sorted(set(dc.values()))
    reds = [c for c in colorset if c.startswith("#") and c[1] in "CFEDBA"
            and c not in ("#318E72", "#216A5B")]
    print(f"    {os.path.basename(f)}")
    print(f"        고유 날짜색: {colorset}")
    print(f"        빨강계열 의심: {reds if reds else '없음 (모두 초록/기타)'}")

print()
# 양식 템플릿
yangsik = os.path.join(BASE, "templates", "2026년 6월 식단표(양식).pptx")
analyze(yangsik, "양식 템플릿")
